import streamlit as st
import pandas as pd
import numpy as np
import plotly.express as px
import plotly.graph_objects as go
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import io
import tempfile
import os
import arabic_reshaper
from bidi.algorithm import get_display

# Set page config
st.set_page_config(page_title="إحصائيات التلاميذ", layout="wide")

# Apply RTL (Right-to-Left) styling for Arabic
st.markdown("""
<style>
    /* Main container RTL */
    .main .block-container {
        direction: rtl;
        text-align: right;
    }
    
    /* Sidebar RTL */
    [data-testid="stSidebar"] {
        direction: rtl;
        text-align: right;
    }
    
    /* Headers and text */
    h1, h2, h3, h4, h5, h6, p, span, div, label {
        direction: rtl;
        text-align: right;
    }
    
    /* Metrics */
    [data-testid="stMetric"] {
        direction: rtl;
        text-align: center;
    }
    
    /* DataFrames */
    .stDataFrame {
        direction: rtl;
    }
    
    /* Tabs */
    .stTabs [data-baseweb="tab-list"] {
        direction: rtl;
        justify-content: flex-end;
    }
    
    /* Selectbox and inputs */
    .stSelectbox, .stMultiSelect, .stTextInput {
        direction: rtl;
    }
    
    /* Expander */
    .streamlit-expanderHeader {
        direction: rtl;
        text-align: right;
    }
    
    /* Info boxes */
    .stAlert {
        direction: rtl;
        text-align: right;
    }
    
    /* Download buttons */
    .stDownloadButton {
        direction: rtl;
    }
    
    /* Checkbox */
    .stCheckbox {
        direction: rtl;
    }
    
    /* Make columns flow RTL */
    [data-testid="column"] {
        direction: rtl;
    }
    
    /* File uploader */
    [data-testid="stFileUploader"] {
        direction: rtl;
    }
    
    /* Caption text */
    .stCaption {
        direction: rtl;
        text-align: right;
    }
    
    /* Markdown content */
    .stMarkdown {
        direction: rtl;
        text-align: right;
    }
    
    /* Tables inside dataframes - align text right */
    .dataframe th, .dataframe td {
        text-align: right !important;
    }
</style>
""", unsafe_allow_html=True)

# ============ ARABIC TEXT FIXER ============
def fix_arabic(text):
    """Reshape and reorder Arabic text for correct rendering in Plotly/Charts"""
    if not text or not isinstance(text, str):
        return text
    try:
        # Check if text contains Arabic characters
        if any('\u0600' <= char <= '\u06FF' for char in text):
            reshaped_text = arabic_reshaper.reshape(text)
            return get_display(reshaped_text)
    except Exception:
        pass
    return text


# ============ POWERPOINT GENERATION IMPORTS & HELPERS ============
from pptx.oxml.ns import qn
from pptx.dml.color import RGBColor

# Define color schemes for fancy styling
PRIMARY_COLOR = RGBColor(0, 112, 192)      # Blue
SECONDARY_COLOR = RGBColor(0, 176, 80)    # Green
ACCENT_COLOR = RGBColor(255, 192, 0)      # Gold
DARK_COLOR = RGBColor(44, 62, 80)         # Dark blue-gray
LIGHT_COLOR = RGBColor(236, 240, 241)     # Light gray

def set_rtl(text_frame):
    """Set Right-to-Left direction on text frame for Arabic"""
    try:
        for paragraph in text_frame.paragraphs:
            pPr = paragraph._p.get_or_add_pPr()
            pPr.set(qn('a:rtl'), '1')
    except Exception:
        pass

def set_paragraph_rtl(paragraph):
    """Set Right-to-Left direction on a paragraph"""
    try:
        pPr = paragraph._p.get_or_add_pPr()
        pPr.set(qn('a:rtl'), '1')
    except Exception:
        pass

def add_gradient_background(slide, color1, color2, angle=90):
    """Add gradient background to slide"""
    try:
        background = slide.background
        fill = background.fill
        fill.gradient()
        fill.gradient_angle = angle
        fill.gradient_stops[0].color.rgb = color1
        fill.gradient_stops[1].color.rgb = color2
    except Exception:
        pass

def add_decorative_shape(slide, shape_type, left, top, width, height, color, transparency=0.3):
    """Add decorative shape"""
    try:
        shape = slide.shapes.add_shape(shape_type, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        shape.fill.fore_color.brightness = transparency
    except Exception:
        pass

def add_title_slide(prs, title, subtitle=""):
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    add_gradient_background(slide, RGBColor(25, 55, 95), RGBColor(45, 85, 135))
    add_decorative_shape(slide, MSO_SHAPE.OVAL, Inches(-2), Inches(-2), Inches(6), Inches(6), RGBColor(255, 255, 255), 0.9)
    add_decorative_shape(slide, MSO_SHAPE.OVAL, Inches(10), Inches(4), Inches(5), Inches(5), RGBColor(255, 255, 255), 0.92)
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = ACCENT_COLOR
    top_bar.line.fill.background()
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(12.333), Inches(1.5))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER
    set_paragraph_rtl(title_para)
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
        sub_frame = subtitle_box.text_frame
        sub_para = sub_frame.paragraphs[0]
        sub_para.text = subtitle
        sub_para.font.size = Pt(24)
        sub_para.font.color.rgb = RGBColor(200, 220, 240)
        sub_para.alignment = PP_ALIGN.CENTER
        set_paragraph_rtl(sub_para)
    bottom_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(5.5), Inches(5.333), Inches(0.05))
    bottom_line.fill.solid()
    bottom_line.fill.fore_color.rgb = ACCENT_COLOR
    bottom_line.line.fill.background()
    return slide

def add_content_slide(prs, title, slide_num=None):
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    add_gradient_background(slide, RGBColor(248, 249, 250), RGBColor(233, 236, 239))
    header_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = PRIMARY_COLOR
    header_bar.line.fill.background()
    accent_stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.1), Inches(13.333), Inches(0.08))
    accent_stripe.fill.solid()
    accent_stripe.fill.fore_color.rgb = ACCENT_COLOR
    accent_stripe.line.fill.background()
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.333), Inches(0.8))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.RIGHT
    set_paragraph_rtl(title_para)
    corner_shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, Inches(0), Inches(6), Inches(1.333), Inches(1.5))
    corner_shape.fill.solid()
    corner_shape.fill.fore_color.rgb = PRIMARY_COLOR
    corner_shape.line.fill.background()
    corner_shape.rotation = 270
    if slide_num is not None:
        num_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(12.4), Inches(6.85), Inches(0.6), Inches(0.6))
        num_circle.fill.solid()
        num_circle.fill.fore_color.rgb = PRIMARY_COLOR
        num_circle.line.fill.background()
        num_txt = slide.shapes.add_textbox(Inches(12.4), Inches(6.92), Inches(0.6), Inches(0.45))
        ntf = num_txt.text_frame
        np = ntf.paragraphs[0]
        np.text = str(slide_num)
        np.font.size = Pt(16)
        np.font.bold = True
        np.font.color.rgb = RGBColor(255, 255, 255)
        np.alignment = PP_ALIGN.CENTER
    return slide

def check_kaleido_available():
    try:
        import kaleido
        test_fig = go.Figure()
        test_fig.to_image(format="png", width=100, height=100)
        return True
    except Exception:
        return False

KALEIDO_AVAILABLE = check_kaleido_available()

def fig_to_image(fig):
    if not KALEIDO_AVAILABLE:
        return None
    try:
        img_bytes = fig.to_image(format="png", width=900, height=500, scale=2)
        return io.BytesIO(img_bytes)
    except Exception:
        return None

def add_toc_slide(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_gradient_background(slide, RGBColor(248, 249, 250), RGBColor(233, 236, 239))
    side_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(13.033), Inches(0), Inches(0.3), Inches(7.5))
    side_bar.fill.solid()
    side_bar.fill.fore_color.rgb = PRIMARY_COLOR
    side_bar.line.fill.background()
    title_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.9))
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = PRIMARY_COLOR
    title_bg.line.fill.background()
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "📋 فهرس المحتويات"
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER
    set_paragraph_rtl(title_para)
    toc_items = [
        ("1", "📈 الإحصائيات العامة", PRIMARY_COLOR),
        ("2", "📊 توزيع شرائح المعدلات", SECONDARY_COLOR),
        ("3", "📈 توزيع المعدلات", PRIMARY_COLOR),
        ("4", "📚 متوسط المعدلات حسب المادة", SECONDARY_COLOR),
        ("5", "📊 المخطط الصندوقي", PRIMARY_COLOR),
        ("6", "🏆 أفضل وأضعف التلاميذ", SECONDARY_COLOR),
        ("7", "💡 أهم الملاحظات", PRIMARY_COLOR),
        ("8", "🔬 مقارنة العلوم والآداب", SECONDARY_COLOR),
        ("9", "🎨 مواد التفتح", PRIMARY_COLOR),
        ("10", "🌐 الكفاءة اللغوية", SECONDARY_COLOR),
        ("11", "💡 التوصيات", PRIMARY_COLOR)
    ]
    y_start = 1.5
    for i, (num, text, color) in enumerate(toc_items):
        if i < 6:
            x_pos, y_pos = 7.0, y_start + (i * 0.45)
        else:
            x_pos, y_pos = 0.8, y_start + ((i - 6) * 0.45)
        item_box = slide.shapes.add_textbox(Inches(x_pos), Inches(y_pos + 0.05), Inches(5.0), Inches(0.4))
        item_para = item_box.text_frame.paragraphs[0]
        item_para.text = text
        item_para.font.size = Pt(18)
        item_para.alignment = PP_ALIGN.RIGHT
        set_paragraph_rtl(item_para)
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x_pos + 5.1), Inches(y_pos), Inches(0.4), Inches(0.4))
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        num_box = slide.shapes.add_textbox(Inches(x_pos + 5.1), Inches(y_pos + 0.05), Inches(0.4), Inches(0.35))
        num_para = num_box.text_frame.paragraphs[0]
        num_para.text = num
        num_para.font.size = Pt(14)
        num_para.font.bold = True
        num_para.font.color.rgb = RGBColor(255, 255, 255)
        num_para.alignment = PP_ALIGN.CENTER
    bottom_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), Inches(6.8), Inches(9.333), Inches(0.05))
    bottom_shape.fill.solid()
    bottom_shape.fill.fore_color.rgb = ACCENT_COLOR
    bottom_shape.line.fill.background()
    return slide

def add_stat_card(slide, x, y, width, height, title, value, icon, bg_color, text_color=RGBColor(255,255,255)):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    card.line.fill.background()
    icon_box = slide.shapes.add_textbox(Inches(x + 1), Inches(y + 0.15), Inches(width - 0.2), Inches(0.5))
    icon_para = icon_box.text_frame.paragraphs[0]
    icon_para.text = icon
    icon_para.font.size = Pt(28)
    icon_para.alignment = PP_ALIGN.CENTER
    value_box = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.6), Inches(width - 0.2), Inches(0.6))
    value_para = value_box.text_frame.paragraphs[0]
    value_para.text = str(value)
    value_para.font.size = Pt(32)
    value_para.font.bold = True
    value_para.font.color.rgb = text_color
    value_para.alignment = PP_ALIGN.CENTER
    title_box = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 1.15), Inches(width - 0.2), Inches(0.4))
    title_para = title_box.text_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(14)
    title_para.font.color.rgb = RGBColor(240, 240, 240)
    title_para.alignment = PP_ALIGN.CENTER
    set_paragraph_rtl(title_para)

def add_bracket_card(slide, x, y, width, height, emoji, title, count, pct, bg_color, border_color):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    card.line.color.rgb = border_color
    card.line.width = Pt(2)
    tf = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.1), Inches(width - 0.2), Inches(height - 0.2)).text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = f"{emoji} {title}"
    p1.font.size = Pt(14)
    p1.font.bold = True
    p1.font.color.rgb = border_color
    p1.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = f"{count} تلميذ"
    p2.font.size = Pt(24)
    p2.font.bold = True
    p2.font.color.rgb = RGBColor(50, 50, 50)
    p2.alignment = PP_ALIGN.CENTER
    p3 = tf.add_paragraph()
    p3.text = f"{pct:.1f}%"
    p3.font.size = Pt(18)
    p3.font.color.rgb = border_color
    p3.alignment = PP_ALIGN.CENTER

def add_fancy_stat(slide, x, y, icon, label, value, bg_color, text_color, width=2.95):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(0.65))
    box.fill.solid()
    box.fill.fore_color.rgb = bg_color
    box.line.color.rgb = text_color
    box.line.width = Pt(2)
    p = slide.shapes.add_textbox(Inches(x + 0.05), Inches(y + 0.05), Inches(width - 0.1), Inches(0.3)).text_frame.paragraphs[0]
    p.text = f"{icon} {label}"
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(100, 100, 100)
    p.alignment = PP_ALIGN.CENTER
    vp = slide.shapes.add_textbox(Inches(x + 0.05), Inches(y + 0.32), Inches(width - 0.1), Inches(0.3)).text_frame.paragraphs[0]
    vp.text = value
    vp.font.size = Pt(18)
    vp.font.bold = True
    vp.font.color.rgb = text_color
    vp.alignment = PP_ALIGN.CENTER

def add_quartile_card(slide, x, y, label, value, color):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(1.95), Inches(0.75))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(250, 250, 250)
    card.line.color.rgb = color
    card.line.width = Pt(1.5)
    lp = slide.shapes.add_textbox(Inches(x + 0.05), Inches(y + 0.08), Inches(1.85), Inches(0.3)).text_frame.paragraphs[0]
    lp.text = label
    lp.font.size = Pt(10)
    lp.font.color.rgb = RGBColor(100, 100, 100)
    lp.alignment = PP_ALIGN.CENTER
    vp = slide.shapes.add_textbox(Inches(x + 0.05), Inches(y + 0.38), Inches(1.85), Inches(0.3)).text_frame.paragraphs[0]
    vp.text = f"{value:.2f}"
    vp.font.size = Pt(16)
    vp.font.bold = True
    vp.font.color.rgb = color
    vp.alignment = PP_ALIGN.CENTER

def add_subject_insight(slide, x, y, icon, title, subject_name, value, bg_color, border_color):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(5.3), Inches(0.85))
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    card.line.color.rgb = border_color
    card.line.width = Pt(2)
    p = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.08), Inches(5.1), Inches(0.35)).text_frame.paragraphs[0]
    p.text = f"{icon} {title}"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = border_color
    p.alignment = PP_ALIGN.RIGHT
    set_paragraph_rtl(p)
    vp = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.43), Inches(5.1), Inches(0.35)).text_frame.paragraphs[0]
    vp.text = f"{subject_name}: {value}"
    vp.font.size = Pt(13)
    vp.font.color.rgb = RGBColor(55, 65, 81)
    vp.alignment = PP_ALIGN.RIGHT
    set_paragraph_rtl(vp)

def add_corr_card(slide, x, y, w, h, icon, title, line1, line2, bg_color, border_color):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    card.line.color.rgb = border_color
    card.line.width = Pt(2)
    p = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.08), Inches(w - 0.2), Inches(0.35)).text_frame.paragraphs[0]
    p.text = f"{icon} {title}"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = border_color
    p.alignment = PP_ALIGN.RIGHT
    set_paragraph_rtl(p)
    p1 = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.4), Inches(w - 0.2), Inches(0.3)).text_frame.paragraphs[0]
    p1.text = line1
    p1.font.size = Pt(11)
    p1.font.color.rgb = RGBColor(55, 65, 81)
    p1.alignment = PP_ALIGN.RIGHT
    set_paragraph_rtl(p1)
    p2 = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.65), Inches(w - 0.2), Inches(0.3)).text_frame.paragraphs[0]
    p2.text = line2
    p2.font.size = Pt(12)
    p2.font.bold = True
    p2.font.color.rgb = border_color
    p2.alignment = PP_ALIGN.RIGHT
    set_paragraph_rtl(p2)

def add_insight_card(slide, x, y, icon, title, subject, value, bg_color, border_color):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(4.3), Inches(0.85))
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    card.line.color.rgb = border_color
    card.line.width = Pt(2)
    p = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.1), Inches(4.1), Inches(0.35)).text_frame.paragraphs[0]
    p.text = f"{icon} {title}"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = border_color
    p.alignment = PP_ALIGN.RIGHT
    set_paragraph_rtl(p)
    val_txt = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.45), Inches(4.1), Inches(0.35)).text_frame.paragraphs[0]
    val_txt.text = f"{subject}: {value}"
    val_txt.font.size = Pt(13)
    val_txt.font.color.rgb = RGBColor(55, 65, 81)
    val_txt.alignment = PP_ALIGN.RIGHT
    set_paragraph_rtl(val_txt)

def generate_slides_for_data(prs, data_df, subject_columns, selected_classes_ppt, title_suffix=""):
    # Title slide
    if len(selected_classes_ppt) == 1:
        classes_text = selected_classes_ppt[0]
    elif len(selected_classes_ppt) <= 3:
        classes_text = ', '.join(selected_classes_ppt)
    else:
        classes_text = f"{len(selected_classes_ppt)} فصول"
    
    add_title_slide(prs, f"📊 إحصائيات نتائج التلاميذ {title_suffix}".strip(), 
                   f"الفصول: {classes_text} | عدد التلاميذ: {len(data_df)}")
    
    # Table of Contents
    add_toc_slide(prs)
    
    # Overall Statistics - Dashboard Style
    slide = add_content_slide(prs, "📈 الإحصائيات العامة", 1)
    
    # Calculate statistics
    total_students = len(data_df)
    avg_grade = data_df['المعدل'].mean()
    max_grade = data_df['المعدل'].max()
    min_grade = data_df['المعدل'].min()
    std_grade = data_df['المعدل'].std()
    num_classes = len(selected_classes_ppt)
    pass_rate = (data_df['المعدل'] >= 10).mean() * 100
    excellent_rate = (data_df['المعدل'] >= 14).mean() * 100
    
    # Row 1: Main metrics (4 cards)
    add_stat_card(slide, 9.8, 1.3, 2.8, 1.6, "عدد التلاميذ", f"{total_students}", "👥", RGBColor(52, 73, 94))
    add_stat_card(slide, 6.8, 1.3, 2.8, 1.6, "المعدل العام", f"{avg_grade:.2f}", "📊", RGBColor(41, 128, 185))
    add_stat_card(slide, 3.8, 1.3, 2.8, 1.6, "نسبة النجاح", f"{pass_rate:.1f}%", "✅", RGBColor(39, 174, 96))
    add_stat_card(slide, 0.8, 1.3, 2.8, 1.6, "عدد الفصول", f"{num_classes}", "🏫", RGBColor(142, 68, 173))
    
    # Row 2: Secondary metrics (4 cards)
    add_stat_card(slide, 9.8, 3.1, 2.8, 1.6, "أعلى معدل", f"{max_grade:.2f}", "🏆", RGBColor(230, 126, 34))
    add_stat_card(slide, 6.8, 3.1, 2.8, 1.6, "أدنى معدل", f"{min_grade:.2f}", "📉", RGBColor(231, 76, 60))
    add_stat_card(slide, 3.8, 3.1, 2.8, 1.6, "الانحراف المعياري", f"{std_grade:.2f}", "📈", RGBColor(52, 152, 219))
    add_stat_card(slide, 0.8, 3.1, 2.8, 1.6, "نسبة التميز (≥14)", f"{excellent_rate:.1f}%", "⭐", RGBColor(241, 196, 15))
    
    # Bottom summary text with explanation
    summary_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.9), Inches(12.3), Inches(1.8))
    summary_frame = summary_box.text_frame
    summary_frame.word_wrap = True
    
    # Performance assessment
    if pass_rate >= 80:
        assessment = "🌟 أداء ممتاز - نسبة نجاح عالية"
        assessment_color = RGBColor(39, 174, 96)
    elif pass_rate >= 60:
        assessment = "✅ أداء جيد - مع إمكانية التحسين"
        assessment_color = RGBColor(241, 196, 15)
    else:
        assessment = "⚠️ يحتاج اهتماماً - نسبة النجاح منخفضة"
        assessment_color = RGBColor(231, 76, 60)
    
    p = summary_frame.paragraphs[0]
    p.text = assessment
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = assessment_color
    p.alignment = PP_ALIGN.CENTER
    set_paragraph_rtl(p)
    
    # Add explanation line
    p2 = summary_frame.add_paragraph()
    p2.text = "📌 نسبة النجاح: معدل ≥ 10 | نسبة التميز: معدل ≥ 14 (جيد جداً/ممتاز)"
    p2.font.size = Pt(14)
    p2.font.color.rgb = RGBColor(100, 100, 100)
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(8)
    set_paragraph_rtl(p2)
    
    # Grade Brackets - Slide 2
    slide = add_content_slide(prs, "📊 توزيع شرائح المعدلات", 2)
    
    below_avg_count = len(data_df[data_df['المعدل'] < 10])
    avg_count = len(data_df[(data_df['المعدل'] >= 10) & (data_df['المعدل'] < 12)])
    good_count = len(data_df[data_df['المعدل'] >= 12])
    total = len(data_df)
    
    add_bracket_card(slide, 9.5, 1.3, 3.2, 1.4, "🔴", "دون المعدل (0-9.99)", below_avg_count, below_avg_count/total*100, 
                    RGBColor(255, 235, 235), RGBColor(231, 76, 60))
    add_bracket_card(slide, 9.5, 2.85, 3.2, 1.4, "🟡", "متوسط (10-11.99)", avg_count, avg_count/total*100,
                    RGBColor(255, 250, 230), RGBColor(241, 196, 15))
    add_bracket_card(slide, 9.5, 4.4, 3.2, 1.4, "🟢", "جيد/ممتاز (12-20)", good_count, good_count/total*100,
                    RGBColor(230, 255, 240), RGBColor(39, 174, 96))
    
    success_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.5), Inches(5.95), Inches(3.2), Inches(0.55))
    success_box.fill.solid()
    success_box.fill.fore_color.rgb = RGBColor(41, 128, 185)
    success_box.line.fill.background()
    
    success_text = slide.shapes.add_textbox(Inches(9.5), Inches(6.0), Inches(3.2), Inches(0.45))
    stf = success_text.text_frame
    sp = stf.paragraphs[0]
    sp.text = f"✅ نسبة النجاح: {(avg_count + good_count)/total*100:.1f}%"
    sp.font.size = Pt(15)
    sp.font.bold = True
    sp.font.color.rgb = RGBColor(255, 255, 255)
    sp.alignment = PP_ALIGN.CENTER
    
    # 3D-style Donut Pie chart
    fig_pie = go.Figure(data=[go.Pie(
        labels=[fix_arabic('دون المعدل<br>(0-9.99)'), fix_arabic('متوسط<br>(10-11.99)'), fix_arabic('جيد/ممتاز<br>(12-20)')],
        values=[below_avg_count, avg_count, good_count],
        hole=0.35,
        marker=dict(colors=['#EF553B', '#FECB52', '#00CC96'], line=dict(color='white', width=3)),
        textinfo='percent+value',
        textfont=dict(size=20, color='white'),
        textposition='inside',
        pull=[0.05, 0.02, 0.02],
        rotation=45,
        direction='clockwise'
    )])
    
    fig_pie.update_layout(
        showlegend=True,
        legend=dict(orientation="h", yanchor="bottom", y=-0.08, xanchor="center", x=0.5, font=dict(size=16)),
        height=680, width=750, margin=dict(t=5, b=40, l=5, r=5), paper_bgcolor='rgba(0,0,0,0)',
        annotations=[dict(text=f'<b>{total}</b><br>{fix_arabic("تلميذ")}', x=0.5, y=0.5, font=dict(size=26, color='#333'), showarrow=False)]
    )
    
    img_stream = fig_to_image(fig_pie)
    if img_stream:
        slide.shapes.add_picture(img_stream, Inches(1.5), Inches(1.0), width=Inches(7.5))
    
    # Grade Distribution Histogram - Slide 3
    slide = add_content_slide(prs, "📈 توزيع المعدلات", 3)
    grades = data_df['المعدل'].dropna()
    grade_mean = grades.mean()
    grade_median = grades.median()
    grade_std = grades.std()
    grade_skew = grades.skew() if len(grades) > 2 else 0
    q1 = grades.quantile(0.25)
    q3 = grades.quantile(0.75)
    iqr = q3 - q1
    passing_rate = (grades >= 10).sum() / len(grades) * 100
    
    if grade_skew > 0.5:
        skew_text, skew_emoji, skew_color = "التوزيع مائل لليمين (معظم الدرجات منخفضة)", "⚠️", RGBColor(231, 76, 60)
    elif grade_skew < -0.5:
        skew_text, skew_emoji, skew_color = "التوزيع مائل لليسار (معظم الدرجات مرتفعة)", "✅", RGBColor(39, 174, 96)
    else:
        skew_text, skew_emoji, skew_color = "التوزيع متماثل تقريباً (طبيعي)", "📊", RGBColor(52, 152, 219)
    
    fig_hist = go.Figure()
    fig_hist.add_trace(go.Histogram(x=grades, nbinsx=20, marker=dict(color='rgba(99, 110, 250, 0.7)', line=dict(color='rgba(99, 110, 250, 1)', width=1))))
    fig_hist.add_vline(x=grade_mean, line_dash="dash", line_color="red", line_width=2, annotation_text=f"{fix_arabic('المتوسط')}: {grade_mean:.2f}", annotation_position="top right")
    fig_hist.add_vline(x=grade_median, line_dash="dot", line_color="green", line_width=2, annotation_text=f"{fix_arabic('الوسيط')}: {grade_median:.2f}", annotation_position="top left")
    fig_hist.add_vline(x=10, line_dash="solid", line_color="orange", line_width=2, annotation_text=fix_arabic("حد النجاح (10)"), annotation_position="bottom right")
    fig_hist.update_layout(height=380, width=580, xaxis_title=fix_arabic("المعدل"), yaxis_title=fix_arabic("عدد التلاميذ"), showlegend=False, margin=dict(t=20, b=40, l=40, r=20))
    
    img_stream = fig_to_image(fig_hist)
    if img_stream:
        slide.shapes.add_picture(img_stream, Inches(0.2), Inches(1.1), width=Inches(5.8))
    
    ititle = slide.shapes.add_textbox(Inches(6.2), Inches(1.1), Inches(6.3), Inches(0.5)).text_frame.paragraphs[0]
    ititle.text = "📊 رؤى إحصائية"
    ititle.font.size = Pt(22)
    ititle.font.bold = True
    ititle.font.color.rgb = PRIMARY_COLOR
    ititle.alignment = PP_ALIGN.RIGHT
    set_paragraph_rtl(ititle)
    
    add_fancy_stat(slide, 6.2, 1.6, "📍", "المتوسط", f"{grade_mean:.2f}", RGBColor(254, 226, 226), RGBColor(220, 38, 38))
    add_fancy_stat(slide, 9.3, 1.6, "📌", "الوسيط", f"{grade_median:.2f}", RGBColor(220, 252, 231), RGBColor(22, 163, 74))
    add_fancy_stat(slide, 6.2, 2.35, "📏", "الانحراف المعياري", f"{grade_std:.2f}", RGBColor(219, 234, 254), RGBColor(37, 99, 235))
    add_fancy_stat(slide, 9.3, 2.35, "📐", "المدى الربيعي", f"{iqr:.2f}", RGBColor(243, 232, 255), RGBColor(147, 51, 234))
    
    pbox = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.2), Inches(3.1), Inches(6.1), Inches(0.7))
    pbox.fill.solid()
    if passing_rate >= 80: pbox.fill.fore_color.rgb, pcolor, picon = RGBColor(220, 252, 231), RGBColor(22, 163, 74), "🏆"
    elif passing_rate >= 60: pbox.fill.fore_color.rgb, pcolor, picon = RGBColor(254, 249, 195), RGBColor(202, 138, 4), "✅"
    else: pbox.fill.fore_color.rgb, pcolor, picon = RGBColor(254, 226, 226), RGBColor(220, 38, 38), "⚠️"
    pbox.line.color.rgb = pcolor
    pbox.line.width = Pt(2)
    pp = slide.shapes.add_textbox(Inches(6.3), Inches(3.2), Inches(5.9), Inches(0.5)).text_frame.paragraphs[0]
    pp.text = f"{picon} نسبة النجاح: {passing_rate:.1f}%"
    pp.font.size = Pt(20)
    pp.font.bold, pp.font.color.rgb, pp.alignment = True, pcolor, PP_ALIGN.CENTER
    
    add_quartile_card(slide, 6.2, 4.3, "الربع الأول (25%)", q1, RGBColor(239, 68, 68))
    add_quartile_card(slide, 8.25, 4.3, "الوسيط (50%)", grade_median, RGBColor(234, 179, 8))
    add_quartile_card(slide, 10.3, 4.3, "الربع الثالث (75%)", q3, RGBColor(34, 197, 94))
    
    ibox = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.2), Inches(5.1), Inches(5.8), Inches(1.1))
    ibox.fill.solid()
    ibox.fill.fore_color.rgb, ibox.line.color.rgb, ibox.line.width = RGBColor(255, 251, 235), skew_color, Pt(2.5)
    ip1 = slide.shapes.add_textbox(Inches(0.3), Inches(5.15), Inches(5.6), Inches(0.35)).text_frame.paragraphs[0]
    ip1.text = "💡 تحليل شكل التوزيع"
    ip1.font.size, ip1.font.bold, ip1.font.color.rgb, ip1.alignment = Pt(14), True, skew_color, PP_ALIGN.RIGHT
    set_paragraph_rtl(ip1)
    itf2 = slide.shapes.add_textbox(Inches(0.3), Inches(5.5), Inches(5.6), Inches(0.6)).text_frame
    itf2.word_wrap = True
    ip2 = itf2.paragraphs[0]
    ip2.text = f"{skew_emoji} {skew_text}"
    ip2.font.size, ip2.font.color.rgb, ip2.alignment = Pt(13), RGBColor(60, 60, 60), PP_ALIGN.RIGHT
    set_paragraph_rtl(ip2)
    ip3 = itf2.add_paragraph()
    ip3.text = f"معامل الالتواء: {grade_skew:.3f}"
    ip3.font.size, ip3.font.color.rgb, ip3.alignment = Pt(11), RGBColor(100, 100, 100), PP_ALIGN.RIGHT
    set_paragraph_rtl(ip3)
    
    # Average by Subject - Slide 4
    slide = add_content_slide(prs, "📚 متوسط المعدلات حسب المادة", 4)
    stats_data_ppt = []
    for col in subject_columns:
        if col in data_df.columns:
            valid_data = data_df[col].dropna()
            if len(valid_data) > 0:
                stats_data_ppt.append({'المادة': col, 'المتوسط': valid_data.mean(), 'الأعلى': valid_data.max(), 'الأقل': valid_data.min(), 'الانحراف المعياري': valid_data.std(), 'عدد الطلاب': len(valid_data), 'نسبة_النجاح': (valid_data >= 10).sum() / len(valid_data) * 100})
    stats_df_ppt = pd.DataFrame(stats_data_ppt)
    stats_df_sorted = stats_df_ppt.sort_values('المتوسط', ascending=True)
    colors = ['#00CC96' if v >= 12 else ('#FECB52' if v >= 10 else '#EF553B') for v in stats_df_sorted['المتوسط']]
    fig_bar = go.Figure(go.Bar(y=[fix_arabic(m) for m in stats_df_sorted['المادة']], x=stats_df_sorted['المتوسط'], orientation='h', marker=dict(color=colors, line=dict(color='white', width=1)), text=[f"{v:.2f}" for v in stats_df_sorted['المتوسط']], textposition='outside'))
    fig_bar.add_vline(x=10, line_dash="dash", line_color="orange", line_width=2, annotation_text=fix_arabic("حد النجاح"), annotation_position="top")
    fig_bar.update_layout(height=420, width=720, xaxis_title=fix_arabic("المتوسط"), yaxis_title="", showlegend=False, margin=dict(t=20, b=40, l=120, r=50), xaxis=dict(range=[0, 20]))
    img_stream = fig_to_image(fig_bar)
    if img_stream: slide.shapes.add_picture(img_stream, Inches(0.2), Inches(1.15), width=Inches(7.2))
    
    if len(stats_df_ppt) > 0:
        best_subject = stats_df_ppt.loc[stats_df_ppt['المتوسط'].idxmax()]
        worst_subject = stats_df_ppt.loc[stats_df_ppt['المتوسط'].idxmin()]
        highest_pass = stats_df_ppt.loc[stats_df_ppt['نسبة_النجاح'].idxmax()]
        lowest_pass = stats_df_ppt.loc[stats_df_ppt['نسبة_النجاح'].idxmin()]
        overall_avg = stats_df_ppt['المتوسط'].mean()
        ititle = slide.shapes.add_textbox(Inches(7.5), Inches(1.15), Inches(5.5), Inches(0.4)).text_frame.paragraphs[0]
        ititle.text = "📊 رؤى تحليلية"
        ititle.font.size, ititle.font.bold, ititle.font.color.rgb, ititle.alignment = Pt(18), True, PRIMARY_COLOR, PP_ALIGN.RIGHT
        set_paragraph_rtl(ititle)
        add_subject_insight(slide, 7.5, 1.6, "🏆", "أفضل مادة (أعلى متوسط)", best_subject['المادة'], f"{best_subject['المتوسط']:.2f}", RGBColor(220, 252, 231), RGBColor(22, 163, 74))
        add_subject_insight(slide, 7.5, 2.55, "⚠️", "أضعف مادة (أدنى متوسط)", worst_subject['المادة'], f"{worst_subject['المتوسط']:.2f}", RGBColor(254, 226, 226), RGBColor(220, 38, 38))
        add_subject_insight(slide, 7.5, 3.5, "✅", "أعلى نسبة نجاح", highest_pass['المادة'], f"{highest_pass['نسبة_النجاح']:.1f}%", RGBColor(219, 234, 254), RGBColor(37, 99, 235))
        add_subject_insight(slide, 7.5, 4.45, "📉", "أدنى نسبة نجاح", lowest_pass['المادة'], f"{lowest_pass['نسبة_النجاح']:.1f}%", RGBColor(254, 249, 195), RGBColor(202, 138, 4))
        abox = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(5.5), Inches(5.3), Inches(0.7))
        abox.fill.solid()
        abox.fill.fore_color.rgb = PRIMARY_COLOR
        ap = slide.shapes.add_textbox(Inches(7.5), Inches(5.6), Inches(5.3), Inches(0.5)).text_frame.paragraphs[0]
        ap.text = f"📈 المتوسط العام لجميع المواد: {overall_avg:.2f}"
        ap.font.size, ap.font.bold, ap.font.color.rgb, ap.alignment = Pt(16), True, RGBColor(255, 255, 255), PP_ALIGN.CENTER
    
    # Subject Failure Analysis - Slide 5
    slide = add_content_slide(prs, "📊 تحليل نسب الرسوب في المواد", 5)
    subject_failure_ppt = []
    for col in subject_columns:
        if col != 'المعدل' and col in data_df.columns:
            sub_data = data_df[col].dropna()
            if len(sub_data) > 0: subject_failure_ppt.append({'المادة': col, 'نسبة الرسوب': (sub_data < 10).mean() * 100})
    if subject_failure_ppt:
        fdf = pd.DataFrame(subject_failure_ppt).sort_values('نسبة الرسوب', ascending=False)
        fig_f = px.bar(fdf, x=[fix_arabic(m) for m in fdf['المادة']], y='نسبة الرسوب', color='نسبة الرسوب', color_continuous_scale='RdYlGn_r', text='نسبة الرسوب')
        fig_f.update_traces(texttemplate='%{text:.1f}%', textposition='outside')
        fig_f.update_layout(height=450, width=1000, xaxis_title=fix_arabic("المادة"), yaxis_title=fix_arabic("نسبة الرسوب"))
        fig_f.add_hline(y=50, line_dash="dash", line_color="red", annotation_text=fix_arabic("خط الخطر"))
        img_stream = fig_to_image(fig_f)
        if img_stream: slide.shapes.add_picture(img_stream, Inches(1.5), Inches(1.3), width=Inches(10))
    
    # Box Plot - Slide 6
    slide = add_content_slide(prs, "📊 توزيع المعدلات حسب المادة (مخطط صندوقي)", 6)
    sdata_list, sstats = [], {}
    for col in subject_columns:
        if col in data_df.columns:
            vd = data_df[col].dropna()
            if len(vd) > 0: sstats[col] = {'median': vd.median(), 'mean': vd.mean(), 'std': vd.std(), 'iqr': vd.quantile(0.75) - vd.quantile(0.25)}
            for g in vd: sdata_list.append({'المادة': col, 'التقدير': g})
    if sdata_list:
        sbdf = pd.DataFrame(sdata_list)
        sbdf['المادة_fixed'] = sbdf['المادة'].apply(fix_arabic)
        fig_b = px.box(sbdf, x='المادة_fixed', y='التقدير', color='المادة_fixed', color_discrete_sequence=px.colors.qualitative.Set2)
        fig_b.update_layout(height=700, width=1200, showlegend=False, xaxis_title=fix_arabic("المادة"), yaxis_title=fix_arabic("التقدير"), font=dict(size=16), margin=dict(t=30, b=60, l=60, r=30))
        img_stream = fig_to_image(fig_b)
        if img_stream: slide.shapes.add_picture(img_stream, Inches(0.3), Inches(1.1), width=Inches(7.8))
        if sstats:
            best_s = max(sstats.items(), key=lambda x: x[1]['median'])
            worst_s = min(sstats.items(), key=lambda x: x[1]['median'])
            most_v = max(sstats.items(), key=lambda x: x[1]['std'])
            most_c = min(sstats.items(), key=lambda x: x[1]['std'])
            ititle = slide.shapes.add_textbox(Inches(8.3), Inches(1.1), Inches(4.5), Inches(0.4)).text_frame.paragraphs[0]
            ititle.text = "💡 رؤى تحليلية"
            ititle.font.size, ititle.font.bold, ititle.font.color.rgb, ititle.alignment = Pt(18), True, PRIMARY_COLOR, PP_ALIGN.RIGHT
            set_paragraph_rtl(ititle)
            add_insight_card(slide, 8.3, 1.5, "🏆", "أفضل مادة (أعلى وسيط)", best_s[0], f"{best_s[1]['median']:.2f}", RGBColor(220, 252, 231), RGBColor(22, 163, 74))
            add_insight_card(slide, 8.3, 2.4, "⚠️", "أضعف مادة (أدنى وسيط)", worst_s[0], f"{worst_s[1]['median']:.2f}", RGBColor(254, 226, 226), RGBColor(220, 38, 38))
            add_insight_card(slide, 8.3, 3.3, "📊", "أكثر تفاوتاً (أعلى انحراف)", most_v[0], f"σ = {most_v[1]['std']:.2f}", RGBColor(254, 249, 195), RGBColor(202, 138, 4))
            add_insight_card(slide, 8.3, 4.2, "✅", "أكثر اتساقاً (أدنى انحراف)", most_c[0], f"σ = {most_c[1]['std']:.2f}", RGBColor(219, 234, 254), RGBColor(37, 99, 235))
    
    # Top & Bottom Performers - Slide 7
    slide = add_content_slide(prs, "🏆 أفضل وأضعف التلاميذ", 7)
    top_10 = data_df[['اسم التلميذ', 'المعدل']].dropna().nlargest(10, 'المعدل')
    bottom_10 = data_df[['اسم التلميذ', 'المعدل']].dropna().nsmallest(10, 'المعدل')
    t_text = "🥇 أفضل 10 تلاميذ:\n"
    emojis = ['🥇', '🥈', '🥉', '4️⃣', '5️⃣', '6️⃣', '7️⃣', '8️⃣', '9️⃣', '🔟']
    for i, (_, r) in enumerate(top_10.iterrows()): t_text += f"{emojis[i]} {r['اسم التلميذ']}: {r['المعدل']:.2f}\n"
    tf = slide.shapes.add_textbox(Inches(6.5), Inches(1.3), Inches(6.3), Inches(5.5)).text_frame
    for line in t_text.strip().split('\n'):
        p = tf.add_paragraph()
        p.text = line
        p.font.size, p.alignment = Pt(16), PP_ALIGN.RIGHT
        set_paragraph_rtl(p)
    b_text = "📉 أضعف 10 تلاميذ (يحتاجون دعماً):\n"
    for _, r in bottom_10.iterrows(): b_text += f"• {r['اسم التلميذ']}: {r['المعدل']:.2f}\n"
    bf = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(6), Inches(5.5)).text_frame
    for line in b_text.strip().split('\n'):
        p = bf.add_paragraph()
        p.text = line
        p.font.size, p.alignment = Pt(16), PP_ALIGN.RIGHT
        set_paragraph_rtl(p)

    # Subject Insights - Slide 8
    slide = add_content_slide(prs, "💡 أهم الملاحظات", 8)
    if len(stats_df_ppt) > 0:
        best_subject = stats_df_ppt.loc[stats_df_ppt['المتوسط'].idxmax()]
        worst_subject = stats_df_ppt.loc[stats_df_ppt['المتوسط'].idxmin()]
        most_consistent = stats_df_ppt.loc[stats_df_ppt['الانحراف المعياري'].idxmin()]
        most_varied = stats_df_ppt.loc[stats_df_ppt['الانحراف المعياري'].idxmax()]
        
        insights_text = f"""
✅ أفضل مادة أداءً: {best_subject['المادة']} (المتوسط: {best_subject['المتوسط']:.2f})
⚠️ مادة تحتاج اهتماماً: {worst_subject['المادة']} (المتوسط: {worst_subject['المتوسط']:.2f})
📊 المادة الأكثر استقراراً: {most_consistent['المادة']} (الانحراف المعياري: {most_consistent['الانحراف المعياري']:.2f})
📈 المادة الأكثر تبايناً: {most_varied['المادة']} (الانحراف المعياري: {most_varied['الانحراف المعياري']:.2f})
🎯 نسبة النجاح الإجمالية: {(avg_count + good_count)/total*100:.1f}%
🌟 نسبة التميز (≥12): {good_count/total*100:.1f}%
        """
        insights_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5))
        insights_frame = insights_box.text_frame
        insights_frame.word_wrap = True
        for line in insights_text.strip().split('\n'):
            p = insights_frame.add_paragraph()
            p.text = line.strip()
            p.font.size = Pt(24)
            p.space_after = Pt(12)
            p.alignment = PP_ALIGN.RIGHT
            set_paragraph_rtl(p)

    # Science vs Humanities Slide - Slide 9
    slide = add_content_slide(prs, "🔬📚 مقارنة العلوم والآداب", 9)
    science_subjects_ppt = ['الرياضيات', 'علوم الحياة والأرض', 'الفيزياء والكيمياء']
    humanities_subjects_ppt = ['اللغة العربية', 'اللغة الفرنسية', 'اللغة الإنجليزية', 'الاجتماعيات', 'التربية الإسلامية']
    science_scores_ppt = []
    humanities_scores_ppt = []
    for col in science_subjects_ppt:
        if col in data_df.columns:
            science_scores_ppt.extend(data_df[col].dropna().tolist())
    for col in humanities_subjects_ppt:
        if col in data_df.columns:
            humanities_scores_ppt.extend(data_df[col].dropna().tolist())
    science_avg_ppt = np.mean(science_scores_ppt) if science_scores_ppt else 0
    humanities_avg_ppt = np.mean(humanities_scores_ppt) if humanities_scores_ppt else 0
    diff_ppt = science_avg_ppt - humanities_avg_ppt
    orientation = "توجه علمي" if diff_ppt > 0.5 else ("توجه أدبي" if diff_ppt < -0.5 else "متوازن")
    sci_hum_text = f"🔬 متوسط المواد العلمية: {science_avg_ppt:.2f}\n(الرياضيات، علوم الحياة والأرض، الفيزياء والكيمياء)\n\n📚 متوسط المواد الأدبية: {humanities_avg_ppt:.2f}\n(العربية، الفرنسية، الإنجليزية، الاجتماعيات، التربية الإسلامية)\n\n📊 الفرق: {diff_ppt:.2f} نقطة\n\n🎯 التوجه العام: {orientation}"
    sci_hum_box = slide.shapes.add_textbox(Inches(6.5), Inches(1.3), Inches(6.3), Inches(5))
    sci_hum_frame = sci_hum_box.text_frame
    sci_hum_frame.word_wrap = True
    for line in sci_hum_text.strip().split('\n'):
        p = sci_hum_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(22)
        p.space_after = Pt(8)
        p.alignment = PP_ALIGN.RIGHT
        set_paragraph_rtl(p)
    comparison_df_ppt = pd.DataFrame({fix_arabic('المجال'): [fix_arabic('المواد العلمية'), fix_arabic('المواد الأدبية')], fix_arabic('المتوسط'): [science_avg_ppt, humanities_avg_ppt]})
    fig_comparison = px.bar(comparison_df_ppt, x=fix_arabic('المجال'), y=fix_arabic('المتوسط'), color=fix_arabic('المجال'), 
                           color_discrete_map={fix_arabic('المواد العلمية'): '#636EFA', fix_arabic('المواد الأدبية'): '#EF553B'}, text=fix_arabic('المتوسط'))
    fig_comparison.update_traces(texttemplate='%{text:.2f}', textposition='outside')
    fig_comparison.update_layout(height=400, width=500, showlegend=False)
    fig_comparison.add_hline(y=10, line_dash="dash", line_color="green")
    img_stream = fig_to_image(fig_comparison)
    if img_stream:
        slide.shapes.add_picture(img_stream, Inches(0.5), Inches(1.3), width=Inches(6))

    # ====== ENRICHMENT SUBJECTS SLIDE ====== Slide 10
    slide = add_content_slide(prs, "🎨 مواد التفتح (الأنشطة)", 10)
    enrichment_subjects_ppt = ['التربية البدنية', 'المعلوميات', 'التربية التشكيلية']
    enrichment_data_ppt = []
    for subj in enrichment_subjects_ppt:
        if subj in data_df.columns:
            # Safety: Ensure data is numeric
            s_data = pd.to_numeric(data_df[subj].astype(str).str.replace(',', '.'), errors='coerce').dropna()
            if len(s_data) > 0:
                avg_val = s_data.mean()
                pass_rate = (s_data >= 10).mean() * 100
                enrichment_data_ppt.append({'المادة': subj, 'المتوسط': avg_val, 'نسبة النجاح': pass_rate})
    if enrichment_data_ppt:
        enrichment_df_ppt = pd.DataFrame(enrichment_data_ppt)
        enrichment_text = "📊 أداء التلاميذ في مواد التفتح:\n\n"
        for _, row in enrichment_df_ppt.iterrows():
            emoji = "✅" if row['المتوسط'] >= 10 else "⚠️"
            enrichment_text += f"{emoji} {row['المادة']}: {row['المتوسط']:.2f} (نجاح: {row['نسبة النجاح']:.0f}%)\n"
        enr_box = slide.shapes.add_textbox(Inches(6.5), Inches(1.3), Inches(6.3), Inches(5))
        enr_frame = enr_box.text_frame
        enr_frame.word_wrap = True
        for line in enrichment_text.strip().split('\n'):
            p = enr_frame.add_paragraph()
            p.text = line
            p.font.size = Pt(20)
            p.space_after = Pt(6)
            p.alignment = PP_ALIGN.RIGHT
            set_paragraph_rtl(p)
        fig_enr = px.bar(enrichment_df_ppt, x=[fix_arabic(m) for m in enrichment_df_ppt['المادة']], y='المتوسط', color='المتوسط', color_continuous_scale='RdYlGn', text='المتوسط')
        fig_enr.update_traces(texttemplate='%{text:.2f}', textposition='outside')
        fig_enr.update_layout(height=500, width=650, showlegend=False)
        fig_enr.add_hline(y=10, line_dash="dash", line_color="green")
        img_stream = fig_to_image(fig_enr)
        if img_stream:
            slide.shapes.add_picture(img_stream, Inches(0.3), Inches(1.2), width=Inches(7))

    # ====== LANGUAGE SUCCESS RATES SLIDE ====== Slide 11
    slide = add_content_slide(prs, "📊 نسبة النجاح في اللغات", 11)
    ar_pass_ppt = (data_df['اللغة العربية'].dropna() >= 10).mean() * 100 if 'اللغة العربية' in data_df.columns else 0
    fr_pass_ppt = (data_df['اللغة الفرنسية'].dropna() >= 10).mean() * 100 if 'اللغة الفرنسية' in data_df.columns else 0
    en_pass_ppt = (data_df['اللغة الإنجليزية'].dropna() >= 10).mean() * 100 if 'اللغة الإنجليزية' in data_df.columns else 0
    pass_df_ppt = pd.DataFrame({fix_arabic('اللغة'): [fix_arabic('العربية'), fix_arabic('الفرنسية'), fix_arabic('الإنجليزية')], fix_arabic('نسبة النجاح %'): [ar_pass_ppt, fr_pass_ppt, en_pass_ppt]})
    fig_pass = px.bar(pass_df_ppt, x=fix_arabic('اللغة'), y=fix_arabic('نسبة النجاح %'), color=fix_arabic('نسبة النجاح %'), color_continuous_scale='RdYlGn', text=fix_arabic('نسبة النجاح %'))
    fig_pass.update_traces(texttemplate='%{text:.1f}%', textposition='outside')
    fig_pass.update_layout(height=400, width=500, title=fix_arabic("نسبة النجاح في كل لغة (≥10)"))
    img_stream_pass = fig_to_image(fig_pass)
    if img_stream_pass:
        slide.shapes.add_picture(img_stream_pass, Inches(0.5), Inches(1.3), width=Inches(6))
    success_analysis = f"📈 نسب النجاح في اللغات:\n\n🇲🇦 العربية: {ar_pass_ppt:.1f}%\n🇫🇷 الفرنسية: {fr_pass_ppt:.1f}%\n🇬🇧 الإنجليزية: {en_pass_ppt:.1f}%\n\n"
    struggling_langs_ppt = []
    if fr_pass_ppt < 50: struggling_langs_ppt.append("الفرنسية")
    if en_pass_ppt < 50: struggling_langs_ppt.append("الإنجليزية")
    if struggling_langs_ppt: success_analysis += f"⚠️ لغات تحتاج دعم: {', '.join(struggling_langs_ppt)}"
    else: success_analysis += "✅ أداء جيد في جميع اللغات"
    success_box = slide.shapes.add_textbox(Inches(6.5), Inches(1.3), Inches(6.3), Inches(5))
    success_frame = success_box.text_frame
    success_frame.word_wrap = True
    for line in success_analysis.strip().split('\n'):
        p = success_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(24)
        p.space_after = Pt(10)
        p.alignment = PP_ALIGN.RIGHT
        set_paragraph_rtl(p)
    
    # Language Proficiency Gap Slide - Slide 12
    slide = add_content_slide(prs, "🌐 فجوة الكفاءة اللغوية", 12)
    arabic_avg_ppt = data_df['اللغة العربية'].dropna().mean() if 'اللغة العربية' in data_df.columns else 0
    french_avg_ppt = data_df['اللغة الفرنسية'].dropna().mean() if 'اللغة الفرنسية' in data_df.columns else 0
    english_avg_ppt = data_df['اللغة الإنجليزية'].dropna().mean() if 'اللغة الإنجليزية' in data_df.columns else 0
    foreign_avg_ppt = np.mean([french_avg_ppt, english_avg_ppt]) if (french_avg_ppt > 0 or english_avg_ppt > 0) else 0
    proficiency_gap_ppt = arabic_avg_ppt - foreign_avg_ppt
    lang_text = f"🇲🇦 اللغة العربية (اللغة الأم): {arabic_avg_ppt:.2f}\n🇫🇷 اللغة الفرنسية: {french_avg_ppt:.2f}\n🇬🇧 الإنجليزية: {english_avg_ppt:.2f}\n📊 فجوة الكفاءة (العربية - الأجنبية): {proficiency_gap_ppt:.2f}"
    if proficiency_gap_ppt > 2: lang_text += "\n⚠️ فجوة كبيرة: التلاميذ يواجهون صعوبة في اللغات الأجنبية"
    elif proficiency_gap_ppt > 1: lang_text += "\n📊 فجوة متوسطة: يحتاج تعزيز اللغات الأجنبية"
    else: lang_text += "\n✅ فجوة صغيرة: الأداء متقارب بين اللغات"
    lang_box = slide.shapes.add_textbox(Inches(6.5), Inches(1.3), Inches(6.3), Inches(5))
    lang_frame = lang_box.text_frame
    lang_frame.word_wrap = True
    for line in lang_text.strip().split('\n'):
        p = lang_frame.add_paragraph()
        p.text = line
        p.font.size, p.alignment, p.space_after = Pt(22), PP_ALIGN.RIGHT, Pt(8)
        set_paragraph_rtl(p)
    lang_df_ppt = pd.DataFrame({fix_arabic('اللغة'): [fix_arabic('العربية'), fix_arabic('الفرنسية'), fix_arabic('الإنجليزية')], fix_arabic('المتوسط'): [arabic_avg_ppt, french_avg_ppt, english_avg_ppt], fix_arabic('النوع'): [fix_arabic('اللغة الأم'), fix_arabic('لغة أجنبية'), fix_arabic('لغة أجنبية')]})
    fig_lang = px.bar(lang_df_ppt, x=fix_arabic('اللغة'), y=fix_arabic('المتوسط'), color=fix_arabic('النوع'), color_discrete_map={fix_arabic('اللغة الأم'): '#00CC96', fix_arabic('لغة أجنبية'): '#EF553B'}, text=fix_arabic('المتوسط'))
    fig_lang.update_traces(texttemplate='%{text:.2f}', textposition='outside')
    fig_lang.update_layout(height=400, width=500, showlegend=True, xaxis_title=fix_arabic("اللغة"), yaxis_title=fix_arabic("المتوسط"))
    fig_lang.add_hline(y=10, line_dash="dash", line_color="gray")
    img_stream = fig_to_image(fig_lang)
    if img_stream: slide.shapes.add_picture(img_stream, Inches(0.5), Inches(1.3), width=Inches(6))

    # ====== LANGUAGE GAP DISTRIBUTION SLIDE ====== Slide 13
    slide = add_content_slide(prs, "📊 توزيع الفجوة اللغوية", 13)
    student_gap_ppt = []
    for idx, row in data_df.iterrows():
        arabic_score = row.get('اللغة العربية', np.nan)
        foreign_scores = [row.get(col, np.nan) for col in ['اللغة الفرنسية', 'اللغة الإنجليزية'] if col in data_df.columns]
        foreign_scores = [s for s in foreign_scores if pd.notna(s)]
        if pd.notna(arabic_score) and foreign_scores:
            student_gap_ppt.append(arabic_score - (sum(foreign_scores) / len(foreign_scores)))
    if student_gap_ppt:
        valid_gaps_ppt = [g for g in student_gap_ppt if pd.notna(g)]
        if valid_gaps_ppt:
            pos_gap = sum(1 for g in valid_gaps_ppt if g > 1)
            neg_gap = sum(1 for g in valid_gaps_ppt if g < -1)
            balanced = len(valid_gaps_ppt) - pos_gap - neg_gap
            fig_gap_hist = px.histogram(pd.DataFrame({fix_arabic('الفجوة'): valid_gaps_ppt}), x=fix_arabic('الفجوة'), nbins=20, color_discrete_sequence=['#636EFA'])
            fig_gap_hist.add_vline(x=0, line_dash="dash", line_color="red", annotation_text=fix_arabic("توازن"))
            fig_gap_hist.update_layout(title=fix_arabic("توزيع الفجوة اللغوية"), height=400, width=550, xaxis_title=fix_arabic("الفجوة"), yaxis_title=fix_arabic("عدد التلاميذ"))
            img_stream_gap = fig_to_image(fig_gap_hist)
            if img_stream_gap: slide.shapes.add_picture(img_stream_gap, Inches(0.3), Inches(1.3), width=Inches(6.2))
            gap_analysis = f"📊 تحليل الفجوة اللغوية:\n\n📈 أفضل في العربية: {pos_gap} تلميذ ({pos_gap/len(valid_gaps_ppt)*100:.1f}%)\n⚖️ متوازن: {balanced} تلميذ ({balanced/len(valid_gaps_ppt)*100:.1f}%)\n🌍 أفضل في الأجنبية: {neg_gap} تلميذ ({neg_gap/len(valid_gaps_ppt)*100:.1f}%)\n\n"
            avg_gap = sum(valid_gaps_ppt) / len(valid_gaps_ppt)
            if avg_gap > 1: gap_analysis += "⚠️ غالبية التلاميذ يحتاجون دعماً في اللغات الأجنبية"
            elif avg_gap < -1: gap_analysis += "🌟 غالبية التلاميذ متفوقون في اللغات الأجنبية"
            else: gap_analysis += "✅ توزيع متوازن للكفاءة اللغوية"
            gap_box = slide.shapes.add_textbox(Inches(6.5), Inches(1.3), Inches(6.3), Inches(5))
            gap_frame = gap_box.text_frame
            gap_frame.word_wrap = True
            for line in gap_analysis.strip().split('\n'):
                p = gap_frame.add_paragraph()
                p.text = line
                p.font.size, p.alignment, p.space_after = Pt(22), PP_ALIGN.RIGHT, Pt(8)
                set_paragraph_rtl(p)

    # Correlation Analysis Slide - Slide 14
    slide = add_content_slide(prs, "🔗 تحليل الارتباط بين المواد", 14)
    corr_subjects = [col for col in subject_columns if col in data_df.columns and col != 'المعدل']
    corr_data = data_df[corr_subjects].dropna()
    if len(corr_data) > 5 and len(corr_subjects) > 1:
        corr_matrix = corr_data.corr()
        correlations = []
        for i in range(len(corr_subjects)):
            for j in range(i + 1, len(corr_subjects)):
                correlations.append({'المادة 1': corr_subjects[i], 'المادة 2': corr_subjects[j], 'الارتباط': corr_matrix.iloc[i, j]})
        corr_df = pd.DataFrame(correlations).sort_values('الارتباط', ascending=False, key=abs)
        avg_corr = corr_df['الارتباط'].mean()
        strongest = corr_df.iloc[0] if len(corr_df) > 0 else None
        avg_color = RGBColor(22, 163, 74) if avg_corr >= 0.5 else (RGBColor(202, 138, 4) if avg_corr >= 0.3 else RGBColor(220, 38, 38))
        avg_bg = RGBColor(220, 252, 231) if avg_corr >= 0.5 else (RGBColor(254, 249, 195) if avg_corr >= 0.3 else RGBColor(254, 226, 226))
        add_corr_card(slide, 6.7, 1.2, 6, 1.0, "📊", "متوسط الارتباط العام", "قياس العلاقة بين جميع المواد", f"{avg_corr:.2f}", avg_bg, avg_color)
        if strongest is not None:
            add_corr_card(slide, 6.7, 2.3, 6, 1.0, "🔗", "أقوى ارتباط", f"{strongest['المادة 1']} ↔ {strongest['المادة 2']}", f"معامل الارتباط: {strongest['الارتباط']:.2f}", RGBColor(219, 234, 254), RGBColor(37, 99, 235))
        fig_corr = px.imshow(corr_matrix, labels=dict(x=fix_arabic("المادة"), y=fix_arabic("المادة"), color=fix_arabic("الارتباط")), x=[fix_arabic(s) for s in corr_subjects], y=[fix_arabic(s) for s in corr_subjects], color_continuous_scale='RdBu_r', zmin=-1, zmax=1, text_auto='.2f')
        fig_corr.update_layout(height=800, width=900, margin=dict(t=20, b=60, l=60, r=40), font=dict(size=14))
        img_stream = fig_to_image(fig_corr)
        if img_stream: slide.shapes.add_picture(img_stream, Inches(0.2), Inches(1.1), width=Inches(6.3))

    # At-Risk Students Slide - Slide 15
    slide = add_content_slide(prs, "🚨 التلاميذ المعرضين للخطر", 15)
    at_risk = data_df[data_df['المعدل'] < 9]
    borderline = data_df[(data_df['المعدل'] >= 9) & (data_df['المعدل'] < 10)]
    excellent = data_df[data_df['المعدل'] >= data_df['المعدل'].mean() + 1.5 * data_df['المعدل'].std()]
    risk_text = f"🔴 معرضون للخطر (معدل < 9): {len(at_risk)} تلاميذ\n🟡 على الحافة (معدل 9-10): {len(borderline)} تلاميذ\n⭐ متميزون جداً: {len(excellent)} تلاميذ"
    risk_box = slide.shapes.add_textbox(Inches(6.5), Inches(1.3), Inches(6.3), Inches(4))
    risk_frame = risk_box.text_frame
    risk_frame.word_wrap = True
    for line in risk_text.strip().split('\n'):
        p = risk_frame.add_paragraph()
        p.text = line
        p.font.size, p.alignment, p.space_after = Pt(22), PP_ALIGN.RIGHT, Pt(8)
        set_paragraph_rtl(p)
    if len(at_risk) > 0:
        names = at_risk.nsmallest(5, 'المعدل')[['اسم التلميذ', 'المعدل']]
        names_text = "📋 أسماء التلاميذ الأكثر خطراً:\n"
        for _, r in names.iterrows(): names_text += f"• {r['اسم التلميذ']}: {r['المعدل']:.2f}\n"
        nf = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(6), Inches(4)).text_frame
        for line in names_text.strip().split('\n'):
            p = nf.add_paragraph()
            p.text = line
            p.font.size, p.alignment, p.space_after = Pt(20), PP_ALIGN.RIGHT, Pt(6)
            set_paragraph_rtl(p)

    # Final Recommendations Slide - Slide 16
    slide = add_content_slide(prs, "💡 التوصيات والخلاصة", 16)
    rec_text = "📌 التوصيات الرئيسية:\n"
    if len(at_risk) > 0: rec_text += f"🔴 تدخل عاجل: {len(at_risk)} تلاميذ يحتاجون دعماً مكثفاً\n"
    if len(borderline) > 0: rec_text += f"🟡 متابعة دقيقة: {len(borderline)} تلاميذ على حافة الرسوب\n"
    if len(excellent) > 0: rec_text += f"⭐ متميزون: {len(excellent)} تلاميذ يمكنهم المساعدة\n"
    rec_text += f"\n📊 ملخص الأداء:\n• نسبة النجاح: {pass_rate:.1f}%\n• نسبة التميز: {excellent_rate:.1f}%\n• المعدل العام: {data_df['المعدل'].mean():.2f}"
    rec_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12), Inches(5.5))
    rec_frame = rec_box.text_frame
    rec_frame.word_wrap = True
    for line in rec_text.strip().split('\n'):
        p = rec_frame.add_paragraph()
        p.text = line
        p.font.size, p.alignment, p.space_after = Pt(22), PP_ALIGN.RIGHT, Pt(8)
        set_paragraph_rtl(p)

    # Thank You Slide
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_background(s, RGBColor(0, 100, 80), RGBColor(25, 55, 95))
    tp = s.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.333), Inches(1.2)).text_frame.paragraphs[0]
    tp.text = "شكراً لكم!"
    tp.font.size, tp.font.bold, tp.font.color.rgb, tp.alignment = Pt(60), True, RGBColor(255,255,255), PP_ALIGN.CENTER
    set_paragraph_rtl(tp)

# ============ GENDER DETECTION FUNCTION ============
def detect_gender(name):
    """
    Detect gender from Arabic/Moroccan first names.
    Returns: 'M' for male, 'F' for female, 'U' for unknown
    """
    if pd.isna(name) or not isinstance(name, str):
        return 'U'
    
    # Extract first name (first word)
    first_name = name.strip().split()[0] if name.strip() else ''
    first_name = first_name.strip()
    
    # Common Moroccan/Arabic female names
    female_names = {
        # Names ending with typical female suffixes
        'فاطمة', 'عائشة', 'خديجة', 'مريم', 'زينب', 'أمينة', 'حليمة', 'رقية', 'سعاد', 'نادية',
        'سميرة', 'نجاة', 'لطيفة', 'حسناء', 'سناء', 'هناء', 'دعاء', 'آسية', 'سارة', 'ليلى',
        'نعيمة', 'كريمة', 'رحيمة', 'فتيحة', 'صفية', 'سلمى', 'هدى', 'منى', 'سهام', 'إيمان',
        'سكينة', 'حنان', 'غيثة', 'رجاء', 'وفاء', 'صباح', 'نوال', 'سعيدة', 'جميلة', 'نبيلة',
        'عزيزة', 'حفيظة', 'رشيدة', 'مليكة', 'خولة', 'أسماء', 'بشرى', 'إكرام', 'ابتسام', 'انتصار',
        'سمية', 'علية', 'زهرة', 'ياسمين', 'نسرين', 'شيماء', 'إسراء', 'آية', 'مروة', 'هاجر',
        'سلوى', 'لبنى', 'رانيا', 'دنيا', 'نهى', 'نورة', 'نور', 'سهى', 'ندى', 'هبة', 'لينا',
        'ريم', 'رنا', 'دينا', 'منار', 'ملاك', 'جنة', 'روان', 'تسنيم', 'سجى', 'وئام', 'نجوى',
        'توفيقة', 'مباركة', 'عتيقة', 'يامنة', 'فضيلة', 'زهور', 'بديعة', 'نزهة', 'حورية', 'سعدية',
        'فوزية', 'زكية', 'تركية', 'خيرة', 'عيشة', 'للا', 'أم', 'فاطنة', 'رحمة', 'بركة',
        'إلهام', 'سهيلة', 'نسيمة', 'وسيلة', 'جليلة', 'وردة', 'زوليخة', 'حادة', 'فضمة', 'يطو',
        'إيناس', 'أميمة', 'هند', 'علا', 'رباب', 'شروق', 'غادة', 'لمياء', 'مها', 'ميساء'
    }
    
    # Common Moroccan/Arabic male names
    male_names = {
        'محمد', 'أحمد', 'عبدالله', 'عبد', 'علي', 'حسن', 'حسين', 'عمر', 'خالد', 'يوسف',
        'إبراهيم', 'عثمان', 'سعيد', 'كريم', 'رشيد', 'مصطفى', 'إدريس', 'عزيز', 'حميد', 'مراد',
        'جمال', 'كمال', 'نبيل', 'سمير', 'منير', 'بشير', 'نصير', 'زهير', 'طارق', 'فاروق',
        'صلاح', 'جلال', 'بلال', 'عادل', 'فيصل', 'نوفل', 'جواد', 'فؤاد', 'عماد', 'زياد',
        'أيمن', 'أنس', 'إياد', 'رياض', 'عياض', 'معاذ', 'براء', 'همام', 'هشام', 'وسام',
        'ياسر', 'ناصر', 'منصور', 'عاشور', 'نور الدين', 'صلاح الدين', 'عز الدين', 'بدر الدين',
        'عبدالرحمن', 'عبدالكريم', 'عبدالحق', 'عبدالصمد', 'عبدالعزيز', 'عبدالرحيم', 'عبدالغني',
        'الحسن', 'الحسين', 'المهدي', 'المصطفى', 'الطيب', 'المختار', 'المنصور', 'الشريف',
        'أمين', 'رضوان', 'سفيان', 'عدنان', 'سليمان', 'رمضان', 'شعبان', 'مروان', 'عثمان',
        'حمزة', 'طه', 'يحيى', 'زكريا', 'آدم', 'نوح', 'موسى', 'عيسى', 'داوود', 'سليم',
        'وليد', 'ماجد', 'راشد', 'حامد', 'أسامة', 'قاسم', 'باسم', 'هاشم', 'عصام', 'حسام',
        'فهد', 'سعد', 'فارس', 'أيوب', 'يونس', 'إلياس', 'درويش', 'مبارك', 'ميمون', 'لحسن',
        'بوشعيب', 'العربي', 'الطاهر', 'الصديق', 'عبدو', 'حدو', 'بوزيد', 'مولاي', 'سيدي'
    }
    
    # Check exact match first
    if first_name in female_names:
        return 'F'
    if first_name in male_names:
        return 'M'
    
    # Check if name starts with common male prefixes
    male_prefixes = ['عبد', 'أبو', 'بو', 'سي', 'مولاي', 'سيدي']
    for prefix in male_prefixes:
        if first_name.startswith(prefix):
            return 'M'
    
    # Check endings - common female name endings in Arabic
    female_endings = ['ة', 'اء', 'ى', 'ية']
    for ending in female_endings:
        if first_name.endswith(ending) and len(first_name) > 2:
            return 'F'
    
    # Default to unknown
    return 'U'

def get_gender_stats(df, name_column='اسم التلميذ'):
    """Calculate gender statistics from a dataframe."""
    if name_column not in df.columns:
        return {'M': 0, 'F': 0, 'U': len(df)}
    
    df['_gender'] = df[name_column].apply(detect_gender)
    stats = df['_gender'].value_counts().to_dict()
    
    return {
        'M': stats.get('M', 0),
        'F': stats.get('F', 0),
        'U': stats.get('U', 0)
    }

# File uploader in sidebar
st.sidebar.header("📁 تحميل الملف")
uploaded_file = st.sidebar.file_uploader(
    "اختر ملف Excel",
    type=['xlsx', 'xls'],
    help="قم بتحميل ملف Excel يحتوي على بيانات التلاميذ"
)

if uploaded_file is None:
    st.title("📊 إحصائيات نتائج التلاميذ")
    st.markdown("---")
    st.info("👈 الرجاء تحميل ملف Excel من القائمة الجانبية للبدء")
    st.markdown("""
    ### 📋 تعليمات الاستخدام:
    1. اضغط على **Browse files** في القائمة الجانبية
    2. اختر ملف Excel يحتوي على بيانات التلاميذ
    3. انتظر حتى يتم تحميل البيانات
    4. استعرض الإحصائيات والرسوم البيانية
    """)
    st.stop()

# Extract title from filename
app_title = Path(uploaded_file.name).stem.replace('_', ' - ')

# Title and intro
st.title(f"📊 {app_title}")
st.markdown("---")

# Load data
@st.cache_data
def load_data(file_content, file_name):
    xls = pd.ExcelFile(io.BytesIO(file_content))
    sheet_names = xls.sheet_names
    
    # Filter out the first sheet if it's just a summary
    data_sheets = [s for s in sheet_names if s not in ['ExportMoGenNoteCcParMatie']]
    
    all_data = []
    for sheet in data_sheets:
        df = pd.read_excel(io.BytesIO(file_content), sheet_name=sheet, header=7)
        df['الفصل'] = sheet  # Add class name
        all_data.append(df)
    
    return pd.concat(all_data, ignore_index=True)

# Load the data
file_content = uploaded_file.read()
df = load_data(file_content, uploaded_file.name)

# Convert grades from string (with commas) to float
subject_columns = [
    'اللغة العربية', 'اللغة الفرنسية', 'اللغة الإنجليزية',
    'الاجتماعيات', 'الرياضيات', 'علوم الحياة والأرض',
    'الفيزياء والكيمياء', 'التربية الإسلامية', 'التربية البدنية',
    'المعلوميات', 'التربية التشكيلية', 'التربية الموسيقية', 'المعدل'
]

for col in subject_columns:
    if col in df.columns:
        df[col] = pd.to_numeric(df[col].astype(str).str.replace(',', '.'), errors='coerce')

# Sidebar for filtering
st.sidebar.markdown("---")
st.sidebar.header("🔍 خيارات التصفية")
if 'الفصل' in df.columns:
    classes = ['جميع الفصول'] + list(df['الفصل'].unique())
    selected_class = st.sidebar.selectbox("اختر الفصل:", classes)
    if selected_class == 'جميع الفصول':
        df_filtered = df.copy()
    else:
        df_filtered = df[df['الفصل'] == selected_class].copy()
else:
    df_filtered = df.copy()

# Remove rows with NaN in اسم التلميذ
df_filtered = df_filtered.dropna(subset=['اسم التلميذ'])

# Overall Statistics
st.header("📈 الإحصائيات العامة")
col1, col2, col3, col4 = st.columns(4)

with col1:
    st.metric("عدد التلاميذ", len(df_filtered))

with col2:
    avg_grade = df_filtered['المعدل'].mean()
    st.metric("المعدل العام", f"{avg_grade:.2f}")

with col3:
    max_grade = df_filtered['المعدل'].max()
    st.metric("أعلى معدل", f"{max_grade:.2f}")

with col4:
    min_grade = df_filtered['المعدل'].min()
    st.metric("أدنى معدل", f"{min_grade:.2f}")

st.markdown("---")

# Data Overview Table - Top & Bottom Performers
st.header("🏆 أفضل وأضعف التلاميذ")

st.markdown("""
**نظرة سريعة:** جدول يعرض التلاميذ المتفوقين والمتأخرين مع نقاط قوتهم وضعفهم الرئيسية.
""")

# Function to analyze student strengths and weaknesses
def analyze_student(row, subject_cols):
    scores = {}
    for col in subject_cols:
        if col != 'المعدل' and col in row.index and pd.notna(row.get(col)):
            scores[col] = row[col]
    
    if not scores:
        return "—", "—"
    
    sorted_scores = sorted(scores.items(), key=lambda x: x[1], reverse=True)
    
    # Best subject
    best_subj, best_score = sorted_scores[0]
    
    # Worst subject
    worst_subj, worst_score = sorted_scores[-1]
    
    # Generate strength description
    if best_score >= 18:
        strength = f"متميز في {best_subj} ({best_score:.2f})"
    elif best_score >= 15:
        strength = f"قوي في {best_subj} ({best_score:.2f})"
    else:
        strength = f"أفضل مادة: {best_subj} ({best_score:.2f})"
    
    # Check if struggling
    if worst_score < 10:
        strength += f" | يعاني في {worst_subj} ({worst_score:.2f})"
    
    return strength, worst_subj

# Get subject columns for analysis
analysis_subject_cols = [col for col in subject_columns if col in df_filtered.columns and col != 'المعدل']

# Create top performers table
st.markdown("### 🥇 أفضل التلاميذ")

top_students = df_filtered.nlargest(5, 'المعدل')[['ر.ت', 'اسم التلميذ', 'المعدل'] + analysis_subject_cols].copy()
top_students = top_students.loc[:, ~top_students.columns.duplicated()]  # Remove duplicate columns
top_students['الترتيب'] = range(1, len(top_students) + 1)
top_students['نقاط القوة'] = top_students.apply(lambda row: analyze_student(row, analysis_subject_cols)[0], axis=1)

# Format rank
rank_labels = {1: '🥇 الأول', 2: '🥈 الثاني', 3: '🥉 الثالث', 4: '4️⃣ الرابع', 5: '5️⃣ الخامس'}
top_students['الترتيب'] = top_students['الترتيب'].map(rank_labels)

top_display = top_students[['الترتيب', 'اسم التلميذ', 'المعدل', 'نقاط القوة']].copy()
top_display.loc[:, 'المعدل_formatted'] = top_display['المعدل'].astype(float).round(2).astype(str)
top_display = top_display[['الترتيب', 'اسم التلميذ', 'المعدل_formatted', 'نقاط القوة']]
top_display.columns = ['الترتيب', 'اسم التلميذ', 'المعدل', 'نقاط القوة']

st.dataframe(top_display, use_container_width=True, hide_index=True)

# Highlight top performer
if len(top_students) > 0:
    top_performer = top_students.iloc[0]
    st.success(f"🏆 **المتفوق الأول:** {top_performer['اسم التلميذ']} بمعدل {top_performer['المعدل']:.2f} - {top_performer['نقاط القوة']}")

# Create bottom performers table
st.markdown("### 📉 أضعف التلاميذ")

bottom_students = df_filtered.nsmallest(5, 'المعدل')[['ر.ت', 'اسم التلميذ', 'المعدل'] + analysis_subject_cols].copy()
bottom_students = bottom_students.loc[:, ~bottom_students.columns.duplicated()]  # Remove duplicate columns
bottom_students['الترتيب'] = range(1, len(bottom_students) + 1)

def get_weakness_details(row, subject_cols):
    scores = {col: row[col] for col in subject_cols if pd.notna(row.get(col))}
    if not scores:
        return "—"
    sorted_scores = sorted(scores.items(), key=lambda x: x[1])
    failing_subjects = [(s, sc) for s, sc in sorted_scores if sc < 10]
    if failing_subjects:
        weakest = failing_subjects[0]
        if len(failing_subjects) > 1:
            return f"ضعيف في {weakest[0]} ({weakest[1]:.2f}) + {len(failing_subjects)-1} مواد أخرى"
        else:
            return f"يحتاج دعماً في {weakest[0]} ({weakest[1]:.2f})"
    else:
        best = sorted_scores[-1]
        return f"أقوى مادة: {best[0]} ({best[1]:.2f})"

bottom_students['التحليل'] = bottom_students.apply(lambda row: get_weakness_details(row, analysis_subject_cols), axis=1)

# Find strength even for weak students
bottom_students['نقطة قوة'] = bottom_students.apply(
    lambda row: max([(col, row[col]) for col in analysis_subject_cols if pd.notna(row.get(col))], 
                   key=lambda x: x[1], default=("—", 0))[0] if any(pd.notna(row.get(col)) for col in analysis_subject_cols) else "—",
    axis=1
)

bottom_display = bottom_students[['الترتيب', 'اسم التلميذ', 'المعدل', 'نقطة قوة', 'التحليل']].copy()
bottom_display.loc[:, 'المعدل_formatted'] = bottom_display['المعدل'].astype(float).round(2).astype(str)
bottom_display = bottom_display[['الترتيب', 'اسم التلميذ', 'المعدل_formatted', 'نقطة قوة', 'التحليل']]
bottom_display.columns = ['الترتيب', 'اسم التلميذ', 'المعدل', 'نقطة قوة', 'التحليل']

st.dataframe(bottom_display, use_container_width=True, hide_index=True)

# Quick action recommendation
if len(bottom_students) > 0:
    worst_performer = df_filtered.loc[df_filtered['المعدل'].idxmin()]
    worst_subjects = {col: worst_performer[col] for col in analysis_subject_cols if pd.notna(worst_performer.get(col)) and worst_performer[col] < 10}
    if worst_subjects:
        critical_subject = min(worst_subjects.items(), key=lambda x: x[1])
        st.warning(f"⚠️ **إجراء مقترح:** التلميذ(ة) **{worst_performer['اسم التلميذ']}** يحتاج دعماً عاجلاً في **{critical_subject[0]}** ({critical_subject[1]:.2f})")

# Borderline students (close to passing/failing)
st.markdown("### ⚖️ التلاميذ على الحافة (9-11)")

borderline = df_filtered[(df_filtered['المعدل'] >= 9) & (df_filtered['المعدل'] <= 11)].copy()
if len(borderline) > 0:
    borderline = borderline.sort_values('المعدل')[['ر.ت', 'اسم التلميذ', 'المعدل'] + analysis_subject_cols]
    borderline = borderline.loc[:, ~borderline.columns.duplicated()]  # Remove duplicate columns
    
    borderline['الحالة'] = borderline['المعدل'].apply(
        lambda x: '🔴 قريب من الرسوب' if float(x) < 10 else '🟢 ناجح بفارق بسيط'
    )
    
    def get_weakest_subject(row):
        scores = [(col, row[col]) for col in analysis_subject_cols if col in row.index and pd.notna(row.get(col))]
        if scores:
            weakest = min(scores, key=lambda x: x[1])
            return f"{weakest[0]} ({float(weakest[1]):.2f})"
        return "—"
    
    borderline['المادة المؤثرة'] = borderline.apply(get_weakest_subject, axis=1)
    
    borderline_display = borderline[['اسم التلميذ', 'المعدل', 'الحالة', 'المادة المؤثرة']].head(10).copy()
    borderline_display.loc[:, 'المعدل_formatted'] = borderline_display['المعدل'].astype(float).round(2).astype(str)
    borderline_display = borderline_display[['اسم التلميذ', 'المعدل_formatted', 'الحالة', 'المادة المؤثرة']]
    borderline_display.columns = ['اسم التلميذ', 'المعدل', 'الحالة', 'المادة المؤثرة']
    
    st.dataframe(borderline_display, use_container_width=True, hide_index=True)
    
    # Quick insight
    below_10 = len(borderline[borderline['المعدل'] < 10])
    above_10 = len(borderline[borderline['المعدل'] >= 10])
    st.info(f"📊 من بين {len(borderline)} تلميذ على الحافة: **{below_10}** قريبون من الرسوب، **{above_10}** ناجحون بفارق بسيط")
else:
    st.success("✅ لا يوجد تلاميذ على حافة النجاح/الرسوب")

st.markdown("---")

# Grade Brackets Analysis
st.header("📊 تحليل شرائح المعدلات")

# Create grade brackets
def get_bracket(grade):
    if pd.isna(grade):
        return None
    elif grade < 10:
        return "0 - 9.99 (دون المعدل)"
    elif grade < 12:
        return "10 - 11.99 (متوسط)"
    else:
        return "12 - 20 (جيد/ممتاز)"

df_filtered['Bracket'] = df_filtered['المعدل'].apply(get_bracket)

# Calculate bracket statistics
bracket_stats = df_filtered.groupby('Bracket').agg({
    'المعدل': ['count', 'mean', 'min', 'max', 'std']
}).round(2)
bracket_stats.columns = ['Count', 'Mean', 'Min', 'Max', 'Std Dev']
bracket_stats = bracket_stats.reset_index()

# Display metrics for each bracket
col1, col2, col3 = st.columns(3)

below_avg = df_filtered[df_filtered['المعدل'] < 10]
average = df_filtered[(df_filtered['المعدل'] >= 10) & (df_filtered['المعدل'] < 12)]
good = df_filtered[df_filtered['المعدل'] >= 12]

with col1:
    st.markdown("### 🔴 دون المعدل (0 - 9.99)")
    st.metric("عدد التلاميذ", len(below_avg))
    if len(below_avg) > 0:
        st.metric("النسبة المئوية", f"{len(below_avg)/len(df_filtered)*100:.1f}%")
        st.metric("متوسط المعدل", f"{below_avg['المعدل'].mean():.2f}")

with col2:
    st.markdown("### 🟡 متوسط (10 - 11.99)")
    st.metric("عدد التلاميذ", len(average))
    if len(average) > 0:
        st.metric("النسبة المئوية", f"{len(average)/len(df_filtered)*100:.1f}%")
        st.metric("متوسط المعدل", f"{average['المعدل'].mean():.2f}")

with col3:
    st.markdown("### 🟢 جيد/ممتاز (12 - 20)")
    st.metric("عدد التلاميذ", len(good))
    if len(good) > 0:
        st.metric("النسبة المئوية", f"{len(good)/len(df_filtered)*100:.1f}%")
        st.metric("متوسط المعدل", f"{good['المعدل'].mean():.2f}")

# Pie chart for bracket distribution
st.subheader("توزيع المعدلات حسب الشرائح")
bracket_counts = df_filtered['Bracket'].value_counts().reset_index()
bracket_counts.columns = ['Bracket', 'Count']

col1, col2 = st.columns(2)

with col1:
    fig = px.pie(
        bracket_counts,
        values='Count',
        names='Bracket',
        color='Bracket',
        color_discrete_map={
            "0 - 9.99 (دون المعدل)": "#EF553B",
            "10 - 11.99 (متوسط)": "#FECB52",
            "12 - 20 (جيد/ممتاز)": "#00CC96"
        }
    )
    fig.update_traces(textposition='inside', textinfo='percent+value')
    fig.update_layout(height=400)
    st.plotly_chart(fig, use_container_width=True)

with col2:
    # Insights summary
    st.markdown("### 💡 أهم الملاحظات")
    total = len(df_filtered)
    
    # Success rate (>=10)
    success_rate = (len(average) + len(good)) / total * 100 if total > 0 else 0
    st.info(f"**نسبة النجاح (≥10):** {success_rate:.1f}% من التلاميذ ناجحون")
    
    # Excellence rate (>=12)
    excellence_rate = len(good) / total * 100 if total > 0 else 0
    st.success(f"**نسبة التميز (≥12):** {excellence_rate:.1f}% حصلوا على معدل جيد/ممتاز")
    
    # At-risk students
    at_risk_rate = len(below_avg) / total * 100 if total > 0 else 0
    if at_risk_rate > 0:
        st.warning(f"**تلاميذ يحتاجون دعماً (<10):** {at_risk_rate:.1f}% يحتاجون متابعة إضافية")
    
    # Performance summary
    if success_rate >= 80:
        st.markdown("✅ **الأداء العام:** ممتاز - معظم التلاميذ ناجحون")
    elif success_rate >= 60:
        st.markdown("⚠️ **الأداء العام:** جيد - الأغلبية ناجحون مع إمكانية التحسن")
    else:
        st.markdown("🚨 **الأداء العام:** يحتاج اهتماماً - كثير من التلاميذ يواجهون صعوبات")

# Students list by bracket
st.subheader("📋 التلاميذ حسب الشريحة")
bracket_tab1, bracket_tab2, bracket_tab3 = st.tabs(["🔴 دون المعدل", "🟡 متوسط", "🟢 جيد/ممتاز"])

with bracket_tab1:
    if len(below_avg) > 0:
        st.dataframe(below_avg[['اسم التلميذ', 'الفصل', 'المعدل']].sort_values('المعدل', ascending=False), use_container_width=True)
    else:
        st.success("لا يوجد تلاميذ في هذه الشريحة!")

with bracket_tab2:
    if len(average) > 0:
        st.dataframe(average[['اسم التلميذ', 'الفصل', 'المعدل']].sort_values('المعدل', ascending=False), use_container_width=True)
    else:
        st.info("لا يوجد تلاميذ في هذه الشريحة")

with bracket_tab3:
    if len(good) > 0:
        st.dataframe(good[['اسم التلميذ', 'الفصل', 'المعدل']].sort_values('المعدل', ascending=False), use_container_width=True)
    else:
        st.info("لا يوجد تلاميذ في هذه الشريحة")

st.markdown("---")

# Detailed Statistics by Subject
st.header("📚 إحصائيات حسب المادة")

# Calculate statistics for each subject
stats_data = []
for col in subject_columns:
    if col in df_filtered.columns:
        valid_data = df_filtered[col].dropna()
        if len(valid_data) > 0:
            stats_data.append({
                'المادة': col,
                'المتوسط': valid_data.mean(),
                'الأعلى': valid_data.max(),
                'الأقل': valid_data.min(),
                'الانحراف المعياري': valid_data.std(),
                'عدد الطلاب': len(valid_data)
            })

stats_df = pd.DataFrame(stats_data)

# Display table
st.dataframe(
    stats_df.style.format({
        'المتوسط': '{:.2f}',
        'الأعلى': '{:.2f}',
        'الأقل': '{:.2f}',
        'الانحراف المعياري': '{:.2f}'
    }),
    use_container_width=True
)

st.markdown("---")

# Visualizations
st.header("📊 الرسوم البيانية")

col1, col2 = st.columns(2)

# Average grades by subject
with col1:
    st.subheader("متوسط المعدلات حسب المادة")
    fig = px.bar(
        stats_df.sort_values('المتوسط', ascending=True),
        x='المتوسط',
        y='المادة',
        orientation='h',
        color='المتوسط',
        color_continuous_scale='Viridis'
    )
    fig.update_layout(height=400)
    st.plotly_chart(fig, use_container_width=True)

# Grade distribution
with col2:
    st.subheader("توزيع المعدلات")
    fig = px.histogram(
        df_filtered,
        x='المعدل',
        nbins=20,
        color_discrete_sequence=['#636EFA']
    )
    fig.add_vline(df_filtered['المعدل'].mean(), line_dash="dash", line_color="red", 
                   annotation_text=f"المتوسط: {df_filtered['المعدل'].mean():.2f}")
    fig.update_layout(height=400)
    st.plotly_chart(fig, use_container_width=True)

st.markdown("---")

# Student Rankings
st.header("🏆 أفضل 10 تلاميذ حسب المعدل")
top_students = df_filtered[['اسم التلميذ', 'المعدل']].dropna().nlargest(10, 'المعدل')
st.dataframe(top_students.reset_index(drop=True), use_container_width=True)

st.markdown("---")

# Performance by Subject - Box Plot
st.header("📊 توزيع المعدلات حسب المادة")

st.markdown("""
**📖 كيفية قراءة هذا الرسم البياني:**
- **الصندوق** يوضح أين تقع معظم معدلات التلاميذ (50% الوسطى)
- **الخط داخل الصندوق** هو الوسيط (المعدل الأوسط)
- **الشعيرات** (الخطوط الممتدة من الصندوق) توضح نطاق المعدلات النموذجية
- **النقاط خارج** الشعيرات هي قيم شاذة (معدلات مرتفعة أو منخفضة بشكل غير عادي)
- **صندوق أطول** يعني تباين أكبر في المعدلات لتلك المادة
- **صندوق في موضع أعلى** يعني أداء عام أفضل في تلك المادة
""")

subject_data = []
for col in subject_columns:
    if col in df_filtered.columns:
        valid_data = df_filtered[col].dropna()
        for grade in valid_data:
            subject_data.append({'المادة': col, 'التقدير': grade})

if subject_data:
    subject_box_df = pd.DataFrame(subject_data)
    fig = px.box(subject_box_df, x='المادة', y='التقدير', color='المادة')
    fig.update_layout(height=500, showlegend=False)
    st.plotly_chart(fig, use_container_width=True)
    
    # Add subject-specific insights
    st.markdown("### 📈 ملاحظات حول المواد")
    col1, col2 = st.columns(2)
    
    with col1:
        # Best performing subject
        best_subject = stats_df.loc[stats_df['المتوسط'].idxmax()]
        st.success(f"**أفضل مادة أداءً:** {best_subject['المادة']} (المتوسط: {best_subject['المتوسط']:.2f})")
        
        # Most consistent subject (lowest std dev)
        most_consistent = stats_df.loc[stats_df['الانحراف المعياري'].idxmin()]
        st.info(f"**الأكثر استقراراً:** {most_consistent['المادة']} (الانحراف المعياري: {most_consistent['الانحراف المعياري']:.2f})")
    
    with col2:
        # Subject needing attention
        worst_subject = stats_df.loc[stats_df['المتوسط'].idxmin()]
        st.warning(f"**تحتاج اهتماماً:** {worst_subject['المادة']} (المتوسط: {worst_subject['المتوسط']:.2f})")
        
        # Most varied subject (highest std dev)
        most_varied = stats_df.loc[stats_df['الانحراف المعياري'].idxmax()]
        st.info(f"**الأكثر تبايناً:** {most_varied['المادة']} (الانحراف المعياري: {most_varied['الانحراف المعياري']:.2f})")

st.markdown("---")

# Science vs Humanities Analysis
st.header("🔬📚 مقارنة العلوم والآداب")

st.markdown("""
**تحليل توجه الفصل:** هل التلاميذ أفضل في المواد العلمية أم الأدبية؟
""")

# Define subject groups
science_subjects = ['الرياضيات', 'علوم الحياة والأرض', 'الفيزياء والكيمياء']
humanities_subjects = ['اللغة العربية', 'اللغة الفرنسية', 'اللغة الإنجليزية', 'الاجتماعيات', 'التربية الإسلامية']

# Calculate averages for each group
science_scores = []
humanities_scores = []

for col in science_subjects:
    if col in df_filtered.columns:
        valid_data = df_filtered[col].dropna()
        science_scores.extend(valid_data.tolist())

for col in humanities_subjects:
    if col in df_filtered.columns:
        valid_data = df_filtered[col].dropna()
        humanities_scores.extend(valid_data.tolist())

science_avg = np.mean(science_scores) if science_scores else 0
humanities_avg = np.mean(humanities_scores) if humanities_scores else 0

# Per-student comparison
student_science_avg = []
student_humanities_avg = []

for idx, row in df_filtered.iterrows():
    sci_vals = [row[col] for col in science_subjects if col in df_filtered.columns and pd.notna(row.get(col))]
    hum_vals = [row[col] for col in humanities_subjects if col in df_filtered.columns and pd.notna(row.get(col))]
    
    student_science_avg.append(np.mean(sci_vals) if sci_vals else np.nan)
    student_humanities_avg.append(np.mean(hum_vals) if hum_vals else np.nan)

# Display comparison
col1, col2, col3 = st.columns(3)

with col1:
    st.markdown("### 🔬 المواد العلمية")
    st.metric("المتوسط العام", f"{science_avg:.2f}")
    st.caption(f"الرياضيات، علوم الحياة والأرض، الفيزياء والكيمياء")

with col2:
    st.markdown("### 📚 المواد الأدبية")
    st.metric("المتوسط العام", f"{humanities_avg:.2f}")
    st.caption(f"العربية، الفرنسية، الإنجليزية، الاجتماعيات، التربية الإسلامية")

with col3:
    st.markdown("### 📊 الفرق")
    diff = science_avg - humanities_avg
    if diff > 0:
        st.metric("التوجه", f"علمي (+{diff:.2f})", delta=f"+{diff:.2f}")
    elif diff < 0:
        st.metric("التوجه", f"أدبي ({diff:.2f})", delta=f"{diff:.2f}")
    else:
        st.metric("التوجه", "متوازن", delta="0.00")

# Visualization
col1, col2 = st.columns(2)

with col1:
    # Bar chart comparison
    comparison_df = pd.DataFrame({
        'المجال': ['المواد العلمية 🔬', 'المواد الأدبية 📚'],
        'المتوسط': [science_avg, humanities_avg]
    })
    
    fig = px.bar(
        comparison_df,
        x='المجال',
        y='المتوسط',
        color='المجال',
        color_discrete_map={
            'المواد العلمية 🔬': '#636EFA',
            'المواد الأدبية 📚': '#EF553B'
        },
        text='المتوسط'
    )
    fig.update_traces(texttemplate='%{text:.2f}', textposition='outside')
    fig.update_layout(height=400, showlegend=False)
    fig.add_hline(y=10, line_dash="dash", line_color="green", annotation_text="معدل النجاح (10)")
    st.plotly_chart(fig, use_container_width=True)

with col2:
    # Detailed subject comparison
    subject_comparison = []
    for col in science_subjects:
        if col in df_filtered.columns:
            avg = df_filtered[col].dropna().mean()
            subject_comparison.append({'المادة': col, 'المتوسط': avg, 'المجال': 'علمي'})
    
    for col in humanities_subjects:
        if col in df_filtered.columns:
            avg = df_filtered[col].dropna().mean()
            subject_comparison.append({'المادة': col, 'المتوسط': avg, 'المجال': 'أدبي'})
    
    if subject_comparison:
        subject_comp_df = pd.DataFrame(subject_comparison)
        fig = px.bar(
            subject_comp_df.sort_values('المتوسط', ascending=True),
            x='المتوسط',
            y='المادة',
            color='المجال',
            orientation='h',
            color_discrete_map={'علمي': '#636EFA', 'أدبي': '#EF553B'}
        )
        fig.update_layout(height=400)
        st.plotly_chart(fig, use_container_width=True)

# Insights
st.markdown("### 💡 تحليل التوجه")

if abs(diff) < 0.5:
    st.success("✅ **الفصل متوازن:** الأداء متقارب بين المواد العلمية والأدبية.")
elif diff >= 2:
    st.info("🔬 **توجه علمي قوي:** التلاميذ يتفوقون بشكل ملحوظ في المواد العلمية.")
elif diff >= 0.5:
    st.info("🔬 **توجه علمي طفيف:** أداء أفضل قليلاً في المواد العلمية.")
elif diff <= -2:
    st.info("📚 **توجه أدبي قوي:** التلاميذ يتفوقون بشكل ملحوظ في المواد الأدبية.")
else:
    st.info("📚 **توجه أدبي طفيف:** أداء أفضل قليلاً في المواد الأدبية.")

# Student distribution by tilt
if len(student_science_avg) == len(df_filtered) and len(student_humanities_avg) == len(df_filtered):
    df_filtered_copy = df_filtered.copy()
    df_filtered_copy['معدل_العلوم'] = student_science_avg
    df_filtered_copy['معدل_الآداب'] = student_humanities_avg
    df_filtered_copy['الفرق'] = df_filtered_copy['معدل_العلوم'] - df_filtered_copy['معدل_الآداب']
    
    science_tilt = len(df_filtered_copy[df_filtered_copy['الفرق'] > 0.5])
    humanities_tilt = len(df_filtered_copy[df_filtered_copy['الفرق'] < -0.5])
    balanced = len(df_filtered_copy[(df_filtered_copy['الفرق'] >= -0.5) & (df_filtered_copy['الفرق'] <= 0.5)])
    
    col1, col2, col3 = st.columns(3)
    with col1:
        st.metric("🔬 تلاميذ علميون", science_tilt, help="تلاميذ أداؤهم أفضل في العلوم بفارق > 0.5")
    with col2:
        st.metric("⚖️ تلاميذ متوازنون", balanced, help="تلاميذ متقاربون في الأداء")
    with col3:
        st.metric("📚 تلاميذ أدبيون", humanities_tilt, help="تلاميذ أداؤهم أفضل في الآداب بفارق > 0.5")

st.markdown("---")

# Enrichment Subjects Analysis (مواد التفتح)
st.header("🎨 مواد التفتح وعلاقتها بالتوجه")

st.markdown("""
**تحليل مواد التفتح:** هل التلاميذ العلميون أو الأدبيون أفضل في مواد التفتح؟
""")

# Define enrichment subjects
enrichment_subjects = ['التربية البدنية', 'المعلوميات']

# Calculate enrichment average
enrichment_scores = []
for col in enrichment_subjects:
    if col in df_filtered.columns:
        valid_data = df_filtered[col].dropna()
        enrichment_scores.extend(valid_data.tolist())

enrichment_avg = np.mean(enrichment_scores) if enrichment_scores else 0

# Display enrichment subjects overview
col1, col2, col3, col4 = st.columns(4)

with col1:
    st.markdown("### 🎨 مواد التفتح")
    st.metric("المتوسط العام", f"{enrichment_avg:.2f}")
    st.caption("التربية البدنية، المعلوميات")

# Individual enrichment subjects
enrichment_avgs = {}
for i, col_name in enumerate(enrichment_subjects):
    if col_name in df_filtered.columns:
        avg = df_filtered[col_name].dropna().mean()
        enrichment_avgs[col_name] = avg
        with [col2, col3, col4][i]:
            emoji = ['🕌', '🏃', '💻'][i]
            st.metric(f"{emoji} {col_name}", f"{avg:.2f}")

# Analyze enrichment performance by student orientation
if student_science_avg and student_humanities_avg and len(student_science_avg) == len(student_humanities_avg):
    st.markdown("### 📊 أداء مواد التفتح حسب توجه التلميذ")
    
    # Calculate enrichment average for each student
    student_enrichment_avg = []
    for idx, row in df_filtered.iterrows():
        enr_vals = [row[col] for col in enrichment_subjects if col in df_filtered.columns and pd.notna(row.get(col))]
        if enr_vals:
            student_enrichment_avg.append(np.mean(enr_vals))
        else:
            student_enrichment_avg.append(np.nan)
    
    df_filtered_copy['معدل_التفتح'] = student_enrichment_avg[:len(df_filtered)]
    
    # Categorize students
    science_students = df_filtered_copy[df_filtered_copy['الفرق'] > 0.5]
    humanities_students = df_filtered_copy[df_filtered_copy['الفرق'] < -0.5]
    balanced_students = df_filtered_copy[(df_filtered_copy['الفرق'] >= -0.5) & (df_filtered_copy['الفرق'] <= 0.5)]
    
    # Calculate enrichment averages by orientation
    science_enrichment = science_students['معدل_التفتح'].dropna().mean() if len(science_students) > 0 else 0
    humanities_enrichment = humanities_students['معدل_التفتح'].dropna().mean() if len(humanities_students) > 0 else 0
    balanced_enrichment = balanced_students['معدل_التفتح'].dropna().mean() if len(balanced_students) > 0 else 0
    
    # Display comparison
    col1, col2, col3 = st.columns(3)
    
    with col1:
        st.metric(
            "🔬 العلميون في التفتح", 
            f"{science_enrichment:.2f}" if science_enrichment > 0 else "—",
            help=f"معدل مواد التفتح للتلاميذ ذوي التوجه العلمي ({len(science_students)} تلميذ)"
        )
    
    with col2:
        st.metric(
            "⚖️ المتوازنون في التفتح", 
            f"{balanced_enrichment:.2f}" if balanced_enrichment > 0 else "—",
            help=f"معدل مواد التفتح للتلاميذ المتوازنين ({len(balanced_students)} تلميذ)"
        )
    
    with col3:
        st.metric(
            "📚 الأدبيون في التفتح", 
            f"{humanities_enrichment:.2f}" if humanities_enrichment > 0 else "—",
            help=f"معدل مواد التفتح للتلاميذ ذوي التوجه الأدبي ({len(humanities_students)} تلميذ)"
        )
    
    # Visualization
    col1, col2 = st.columns(2)
    
    with col1:
        # Bar chart for enrichment by orientation
        orientation_enrichment_df = pd.DataFrame({
            'التوجه': ['🔬 علميون', '⚖️ متوازنون', '📚 أدبيون'],
            'معدل التفتح': [science_enrichment, balanced_enrichment, humanities_enrichment],
            'عدد التلاميذ': [len(science_students), len(balanced_students), len(humanities_students)]
        })
        
        fig = px.bar(
            orientation_enrichment_df,
            x='التوجه',
            y='معدل التفتح',
            color='التوجه',
            color_discrete_map={
                '🔬 علميون': '#636EFA',
                '⚖️ متوازنون': '#00CC96',
                '📚 أدبيون': '#EF553B'
            },
            text='معدل التفتح'
        )
        fig.update_traces(texttemplate='%{text:.2f}', textposition='outside')
        fig.update_layout(height=400, showlegend=False, title="معدل مواد التفتح حسب التوجه")
        fig.add_hline(y=10, line_dash="dash", line_color="green", annotation_text="معدل النجاح")
        st.plotly_chart(fig, use_container_width=True)
    
    with col2:
        # Detailed enrichment subjects by orientation
        detailed_data = []
        for subj in enrichment_subjects:
            if subj in df_filtered.columns:
                sci_avg = science_students[subj].dropna().mean() if len(science_students) > 0 else 0
                hum_avg = humanities_students[subj].dropna().mean() if len(humanities_students) > 0 else 0
                bal_avg = balanced_students[subj].dropna().mean() if len(balanced_students) > 0 else 0
                
                detailed_data.append({'المادة': subj, 'المعدل': sci_avg, 'التوجه': 'علميون'})
                detailed_data.append({'المادة': subj, 'المعدل': hum_avg, 'التوجه': 'أدبيون'})
                detailed_data.append({'المادة': subj, 'المعدل': bal_avg, 'التوجه': 'متوازنون'})
        
        if detailed_data:
            detailed_df = pd.DataFrame(detailed_data)
            fig = px.bar(
                detailed_df,
                x='المادة',
                y='المعدل',
                color='التوجه',
                barmode='group',
                color_discrete_map={
                    'علميون': '#636EFA',
                    'متوازنون': '#00CC96',
                    'أدبيون': '#EF553B'
                }
            )
            fig.update_layout(height=400, title="تفصيل مواد التفتح حسب التوجه")
            st.plotly_chart(fig, use_container_width=True)
    
    # Insights
    st.markdown("### 💡 استنتاجات مواد التفتح")
    
    # Determine who performs better
    best_in_enrichment = max(
        [('العلميون', science_enrichment), ('المتوازنون', balanced_enrichment), ('الأدبيون', humanities_enrichment)],
        key=lambda x: x[1] if x[1] > 0 else -999
    )
    
    worst_in_enrichment = min(
        [('العلميون', science_enrichment), ('المتوازنون', balanced_enrichment), ('الأدبيون', humanities_enrichment)],
        key=lambda x: x[1] if x[1] > 0 else 999
    )
    
    if best_in_enrichment[1] > 0 and worst_in_enrichment[1] > 0:
        diff_enrichment = best_in_enrichment[1] - worst_in_enrichment[1]
        
        if diff_enrichment < 0.3:
            st.success("✅ **الأداء متقارب:** جميع التلاميذ بمختلف توجهاتهم لديهم أداء متشابه في مواد التفتح.")
        else:
            st.info(f"📊 **{best_in_enrichment[0]}** هم الأفضل في مواد التفتح بمعدل **{best_in_enrichment[1]:.2f}**، متفوقين على {worst_in_enrichment[0]} بفارق **{diff_enrichment:.2f}** نقطة.")
        
        # Individual subject insights
        for subj in enrichment_subjects:
            if subj in df_filtered.columns:
                sci_avg = science_students[subj].dropna().mean() if len(science_students) > 0 else 0
                hum_avg = humanities_students[subj].dropna().mean() if len(humanities_students) > 0 else 0
                
                if sci_avg > 0 and hum_avg > 0:
                    subj_diff = sci_avg - hum_avg
                    if abs(subj_diff) >= 0.5:
                        if subj_diff > 0:
                            st.caption(f"🔬 **{subj}:** العلميون أفضل بفارق {subj_diff:.2f}")
                        else:
                            st.caption(f"📚 **{subj}:** الأدبيون أفضل بفارق {abs(subj_diff):.2f}")

st.markdown("---")

# Language Proficiency Gap Analysis
st.header("🌐 تحليل فجوة الكفاءة اللغوية")

st.markdown("""
**مقارنة الأداء اللغوي:** هل يواجه التلاميذ صعوبة في اللغات الأجنبية مقارنة بلغتهم الأم؟
""")

# Define language subjects
primary_language = 'اللغة العربية'
foreign_languages = ['اللغة الفرنسية', 'اللغة الإنجليزية']

# Calculate averages
arabic_avg = df_filtered[primary_language].dropna().mean() if primary_language in df_filtered.columns else 0
french_avg = df_filtered['اللغة الفرنسية'].dropna().mean() if 'اللغة الفرنسية' in df_filtered.columns else 0
english_avg = df_filtered['اللغة الإنجليزية'].dropna().mean() if 'اللغة الإنجليزية' in df_filtered.columns else 0
foreign_avg = np.mean([french_avg, english_avg]) if french_avg > 0 or english_avg > 0 else 0

# Language proficiency gap
proficiency_gap = arabic_avg - foreign_avg

# Display metrics
col1, col2, col3, col4 = st.columns(4)

with col1:
    st.markdown("### 🇲🇦 اللغة العربية")
    st.metric("المتوسط", f"{arabic_avg:.2f}")
    st.caption("اللغة الأم")

with col2:
    st.markdown("### 🇫🇷 اللغة الفرنسية")
    st.metric("المتوسط", f"{french_avg:.2f}")
    gap_fr = arabic_avg - french_avg
    if gap_fr > 0:
        st.caption(f"فجوة: -{gap_fr:.2f}")
    else:
        st.caption(f"فرق: +{abs(gap_fr):.2f}")

with col3:
    st.markdown("### 🇬🇧 اللغة الإنجليزية")
    st.metric("المتوسط", f"{english_avg:.2f}")
    gap_en = arabic_avg - english_avg
    if gap_en > 0:
        st.caption(f"فجوة: -{gap_en:.2f}")
    else:
        st.caption(f"فرق: +{abs(gap_en):.2f}")

with col4:
    st.markdown("### 📊 فجوة الكفاءة")
    if proficiency_gap > 0:
        st.metric("الفجوة", f"{proficiency_gap:.2f}", delta=f"-{proficiency_gap:.2f}", delta_color="inverse")
    else:
        st.metric("الفجوة", f"{abs(proficiency_gap):.2f}", delta=f"+{abs(proficiency_gap):.2f}")
    st.caption("الفرق بين العربية واللغات الأجنبية")

# Visualization
col1, col2 = st.columns(2)

with col1:
    # Bar chart for language comparison
    lang_df = pd.DataFrame({
        'اللغة': ['🇲🇦 العربية', '🇫🇷 الفرنسية', '🇬🇧 الإنجليزية'],
        'المتوسط': [arabic_avg, french_avg, english_avg],
        'النوع': ['اللغة الأم', 'لغة أجنبية', 'لغة أجنبية']
    })
    
    fig = px.bar(
        lang_df,
        x='اللغة',
        y='المتوسط',
        color='النوع',
        color_discrete_map={
            'اللغة الأم': '#00CC96',
            'لغة أجنبية': '#EF553B'
        },
        text='المتوسط'
    )
    fig.update_traces(texttemplate='%{text:.2f}', textposition='outside')
    fig.update_layout(height=400, showlegend=True, title="مقارنة الأداء اللغوي")
    fig.add_hline(y=10, line_dash="dash", line_color="gray", annotation_text="معدل النجاح")
    st.plotly_chart(fig, use_container_width=True)

with col2:
    # Radar chart for language skills
    categories = ['العربية', 'الفرنسية', 'الإنجليزية']
    
    fig = go.Figure()
    
    fig.add_trace(go.Scatterpolar(
        r=[arabic_avg, french_avg, english_avg],
        theta=categories,
        fill='toself',
        name='المتوسط الفعلي',
        line_color='#636EFA'
    ))
    
    # Add reference line for passing grade
    fig.add_trace(go.Scatterpolar(
        r=[10, 10, 10],
        theta=categories,
        fill='toself',
        name='معدل النجاح',
        line_color='#00CC96',
        opacity=0.3
    ))
    
    fig.update_layout(
        polar=dict(
            radialaxis=dict(
                visible=True,
                range=[0, 20]
            )
        ),
        showlegend=True,
        title="مخطط الكفاءة اللغوية",
        height=400
    )
    st.plotly_chart(fig, use_container_width=True)

# Per-student language gap analysis
st.markdown("### 📈 توزيع الفجوة اللغوية لدى التلاميذ")

student_arabic = []
student_foreign = []
student_gap = []

for idx, row in df_filtered.iterrows():
    ar = row.get(primary_language) if primary_language in df_filtered.columns else np.nan
    fr = row.get('اللغة الفرنسية') if 'اللغة الفرنسية' in df_filtered.columns else np.nan
    en = row.get('اللغة الإنجليزية') if 'اللغة الإنجليزية' in df_filtered.columns else np.nan
    
    if pd.notna(ar):
        student_arabic.append(ar)
        foreign_vals = [v for v in [fr, en] if pd.notna(v)]
        if foreign_vals:
            foreign_mean = np.mean(foreign_vals)
            student_foreign.append(foreign_mean)
            student_gap.append(ar - foreign_mean)
        else:
            student_foreign.append(np.nan)
            student_gap.append(np.nan)
    else:
        student_arabic.append(np.nan)
        student_foreign.append(np.nan)
        student_gap.append(np.nan)

# Categorize students by gap
positive_gap = sum(1 for g in student_gap if pd.notna(g) and g > 1)  # Better in Arabic
small_gap = sum(1 for g in student_gap if pd.notna(g) and -1 <= g <= 1)  # Balanced
negative_gap = sum(1 for g in student_gap if pd.notna(g) and g < -1)  # Better in foreign languages

col1, col2, col3 = st.columns(3)

with col1:
    st.metric(
        "🇲🇦 أفضل في العربية", 
        positive_gap,
        help="تلاميذ أداؤهم في العربية أفضل من اللغات الأجنبية بفارق > 1"
    )

with col2:
    st.metric(
        "⚖️ متوازنون لغوياً", 
        small_gap,
        help="تلاميذ أداؤهم متقارب في جميع اللغات"
    )

with col3:
    st.metric(
        "🌍 أفضل في الأجنبية", 
        negative_gap,
        help="تلاميذ أداؤهم في اللغات الأجنبية أفضل من العربية بفارق > 1"
    )

# Histogram of language gap
if student_gap:
    valid_gaps = [g for g in student_gap if pd.notna(g)]
    if valid_gaps:
        gap_df = pd.DataFrame({'الفجوة اللغوية': valid_gaps})
        fig = px.histogram(
            gap_df,
            x='الفجوة اللغوية',
            nbins=20,
            color_discrete_sequence=['#636EFA']
        )
        fig.add_vline(x=0, line_dash="dash", line_color="red", annotation_text="توازن")
        fig.update_layout(
            title="توزيع الفجوة اللغوية (العربية - اللغات الأجنبية)",
            xaxis_title="الفجوة (قيم موجبة = أفضل في العربية)",
            yaxis_title="عدد التلاميذ",
            height=350
        )
        st.plotly_chart(fig, use_container_width=True)

# French vs English comparison
st.markdown("### 🇫🇷 vs 🇬🇧 مقارنة اللغتين الأجنبيتين")

col1, col2 = st.columns(2)

with col1:
    fr_en_diff = french_avg - english_avg
    if abs(fr_en_diff) < 0.5:
        st.info("⚖️ **أداء متقارب:** التلاميذ لديهم مستوى متشابه في الفرنسية والإنجليزية.")
    elif fr_en_diff > 0:
        st.info(f"🇫🇷 **الفرنسية أفضل:** التلاميذ يتفوقون في الفرنسية بفارق **{fr_en_diff:.2f}** نقطة.")
    else:
        st.info(f"🇬🇧 **الإنجليزية أفضل:** التلاميذ يتفوقون في الإنجليزية بفارق **{abs(fr_en_diff):.2f}** نقطة.")

with col2:
    # Success rates for each language
    if primary_language in df_filtered.columns:
        ar_pass = (df_filtered[primary_language].dropna() >= 10).mean() * 100
    else:
        ar_pass = 0
    
    if 'اللغة الفرنسية' in df_filtered.columns:
        fr_pass = (df_filtered['اللغة الفرنسية'].dropna() >= 10).mean() * 100
    else:
        fr_pass = 0
    
    if 'اللغة الإنجليزية' in df_filtered.columns:
        en_pass = (df_filtered['اللغة الإنجليزية'].dropna() >= 10).mean() * 100
    else:
        en_pass = 0
    
    pass_df = pd.DataFrame({
        'اللغة': ['العربية', 'الفرنسية', 'الإنجليزية'],
        'نسبة النجاح %': [ar_pass, fr_pass, en_pass]
    })
    
    fig = px.bar(
        pass_df,
        x='اللغة',
        y='نسبة النجاح %',
        color='نسبة النجاح %',
        color_continuous_scale='RdYlGn',
        text='نسبة النجاح %'
    )
    fig.update_traces(texttemplate='%{text:.1f}%', textposition='outside')
    fig.update_layout(height=300, title="نسبة النجاح في كل لغة")
    st.plotly_chart(fig, use_container_width=True)

# Insights
st.markdown("### 💡 استنتاجات الكفاءة اللغوية")

if proficiency_gap > 2:
    st.warning(f"⚠️ **فجوة كبيرة:** التلاميذ يواجهون صعوبة واضحة في اللغات الأجنبية مقارنة بالعربية (فجوة: {proficiency_gap:.2f}). يُنصح بتعزيز برامج تعلم اللغات الأجنبية.")
elif proficiency_gap > 1:
    st.info(f"📊 **فجوة متوسطة:** هناك فرق ملحوظ بين الأداء في العربية واللغات الأجنبية (فجوة: {proficiency_gap:.2f}).")
elif proficiency_gap > 0:
    st.success(f"✅ **فجوة صغيرة:** الأداء متقارب نسبياً بين اللغات (فجوة: {proficiency_gap:.2f}).")
else:
    st.success(f"🌟 **تميز في اللغات الأجنبية:** التلاميذ يؤدون بشكل أفضل في اللغات الأجنبية من العربية!")

# Specific recommendations
if french_avg < 10 or english_avg < 10:
    struggling_langs = []
    if french_avg < 10:
        struggling_langs.append(f"الفرنسية ({french_avg:.2f})")
    if english_avg < 10:
        struggling_langs.append(f"الإنجليزية ({english_avg:.2f})")
    st.caption(f"⚠️ المواد التي تحتاج اهتماماً: {', '.join(struggling_langs)}")

st.markdown("---")

# Correlation Analysis
st.header("🔗 تحليل الارتباط بين المواد")

st.markdown("""
**تحليل العلاقات:** هل النجاح في مادة معينة يتنبأ بالنجاح في مادة أخرى؟
- **ارتباط قوي (> 0.7):** العلاقة قوية جداً
- **ارتباط متوسط (0.4-0.7):** العلاقة معتدلة
- **ارتباط ضعيف (< 0.4):** العلاقة ضعيفة
""")

# Get available subjects for correlation
correlation_subjects = [col for col in subject_columns if col in df_filtered.columns and col != 'المعدل']
correlation_data = df_filtered[correlation_subjects].dropna()

if len(correlation_data) > 5 and len(correlation_subjects) > 1:
    # Calculate correlation matrix
    corr_matrix = correlation_data.corr()
    
    # Heatmap visualization
    st.markdown("### 🗺️ خريطة الارتباط الحرارية")
    
    fig = px.imshow(
        corr_matrix,
        labels=dict(x="المادة", y="المادة", color="معامل الارتباط"),
        x=correlation_subjects,
        y=correlation_subjects,
        color_continuous_scale='RdBu_r',
        zmin=-1,
        zmax=1,
        aspect='auto'
    )
    fig.update_layout(
        height=500,
        title="معاملات الارتباط بين المواد الدراسية"
    )
    # Add correlation values as text
    annotations = []
    for i, row in enumerate(corr_matrix.values):
        for j, val in enumerate(row):
            annotations.append(
                dict(
                    x=j,
                    y=i,
                    text=f"{val:.2f}",
                    showarrow=False,
                    font=dict(color='white' if abs(val) > 0.5 else 'black', size=10)
                )
            )
    fig.update_layout(annotations=annotations)
    st.plotly_chart(fig, use_container_width=True)
    
    # Find strongest correlations (excluding self-correlation)
    st.markdown("### 📊 أقوى العلاقات بين المواد")
    
    # Get upper triangle of correlation matrix (to avoid duplicates)
    correlations = []
    for i in range(len(correlation_subjects)):
        for j in range(i + 1, len(correlation_subjects)):
            correlations.append({
                'المادة 1': correlation_subjects[i],
                'المادة 2': correlation_subjects[j],
                'معامل الارتباط': corr_matrix.iloc[i, j]
            })
    
    corr_df = pd.DataFrame(correlations)
    corr_df['قوة الارتباط'] = corr_df['معامل الارتباط'].abs()
    corr_df = corr_df.sort_values('قوة الارتباط', ascending=False)
    
    # Top 5 strongest correlations
    col1, col2 = st.columns(2)
    
    with col1:
        st.markdown("#### 🔝 أقوى 5 ارتباطات")
        top_5 = corr_df.head(5)
        for idx, row in top_5.iterrows():
            corr_val = row['معامل الارتباط']
            if corr_val >= 0.7:
                emoji = "🟢"
                strength = "قوي جداً"
            elif corr_val >= 0.4:
                emoji = "🟡"
                strength = "متوسط"
            elif corr_val >= 0:
                emoji = "🟠"
                strength = "ضعيف"
            else:
                emoji = "🔴"
                strength = "عكسي"
            
            st.markdown(f"{emoji} **{row['المادة 1']}** ↔ **{row['المادة 2']}**: {corr_val:.2f} ({strength})")
    
    with col2:
        st.markdown("#### 📉 أضعف 5 ارتباطات")
        bottom_5 = corr_df.tail(5).iloc[::-1]
        for idx, row in bottom_5.iterrows():
            corr_val = row['معامل الارتباط']
            if abs(corr_val) < 0.2:
                emoji = "⚪"
                strength = "شبه معدوم"
            elif corr_val < 0:
                emoji = "🔴"
                strength = "عكسي"
            else:
                emoji = "🟠"
                strength = "ضعيف"
            
            st.markdown(f"{emoji} **{row['المادة 1']}** ↔ **{row['المادة 2']}**: {corr_val:.2f} ({strength})")
    
    # Subject-specific correlation analysis
    st.markdown("### 🎯 تحليل ارتباط كل مادة")
    
    selected_subject = st.selectbox(
        "اختر مادة لعرض ارتباطاتها:",
        correlation_subjects,
        key="corr_subject_select"
    )
    
    if selected_subject:
        subject_corr = corr_matrix[selected_subject].drop(selected_subject).sort_values(ascending=False)
        
        col1, col2 = st.columns(2)
        
        with col1:
            # Bar chart of correlations
            corr_chart_df = pd.DataFrame({
                'المادة': subject_corr.index,
                'معامل الارتباط': subject_corr.values
            })
            
            fig = px.bar(
                corr_chart_df,
                x='معامل الارتباط',
                y='المادة',
                orientation='h',
                color='معامل الارتباط',
                color_continuous_scale='RdBu_r',
                range_color=[-1, 1],
                text='معامل الارتباط'
            )
            fig.update_traces(texttemplate='%{text:.2f}', textposition='outside')
            fig.update_layout(height=400, title=f"ارتباطات {selected_subject}")
            fig.add_vline(x=0, line_dash="dash", line_color="gray")
            st.plotly_chart(fig, use_container_width=True)
        
        with col2:
            # Interpretation
            st.markdown(f"#### 💡 تفسير ارتباطات {selected_subject}")
            
            strong_positive = subject_corr[subject_corr >= 0.6]
            moderate_positive = subject_corr[(subject_corr >= 0.4) & (subject_corr < 0.6)]
            weak = subject_corr[(subject_corr > -0.4) & (subject_corr < 0.4)]
            negative = subject_corr[subject_corr <= -0.4]
            
            if len(strong_positive) > 0:
                st.success(f"🟢 **ارتباط قوي مع:** {', '.join(strong_positive.index.tolist())}")
                st.caption("التلاميذ الجيدون في هذه المادة غالباً جيدون في المواد المذكورة")
            
            if len(moderate_positive) > 0:
                st.info(f"🟡 **ارتباط متوسط مع:** {', '.join(moderate_positive.index.tolist())}")
            
            if len(negative) > 0:
                st.warning(f"🔴 **ارتباط عكسي مع:** {', '.join(negative.index.tolist())}")
                st.caption("التلاميذ الجيدون في هذه المادة قد يواجهون صعوبة في المواد المذكورة")
    
    # Scatter plot for specific pairs
    st.markdown("### 📈 رسم الانتشار بين مادتين")
    
    col1, col2 = st.columns(2)
    with col1:
        subject_x = st.selectbox("المادة الأولى (المحور الأفقي):", correlation_subjects, key="scatter_x")
    with col2:
        remaining_subjects = [s for s in correlation_subjects if s != subject_x]
        subject_y = st.selectbox("المادة الثانية (المحور العمودي):", remaining_subjects, key="scatter_y")
    
    if subject_x and subject_y:
        scatter_data = df_filtered[[subject_x, subject_y, 'اسم التلميذ']].dropna()
        
        if len(scatter_data) > 0:
            correlation_value = scatter_data[subject_x].corr(scatter_data[subject_y])
            
            fig = px.scatter(
                scatter_data,
                x=subject_x,
                y=subject_y,
                hover_data=['اسم التلميذ'],
                trendline='ols',
                color_discrete_sequence=['#636EFA']
            )
            fig.update_layout(
                height=450,
                title=f"العلاقة بين {subject_x} و {subject_y} (r = {correlation_value:.2f})"
            )
            # Add quadrant lines at passing grade
            fig.add_hline(y=10, line_dash="dash", line_color="green", opacity=0.5)
            fig.add_vline(x=10, line_dash="dash", line_color="green", opacity=0.5)
            
            st.plotly_chart(fig, use_container_width=True)
            
            # Quadrant analysis
            both_pass = len(scatter_data[(scatter_data[subject_x] >= 10) & (scatter_data[subject_y] >= 10)])
            x_only = len(scatter_data[(scatter_data[subject_x] >= 10) & (scatter_data[subject_y] < 10)])
            y_only = len(scatter_data[(scatter_data[subject_x] < 10) & (scatter_data[subject_y] >= 10)])
            both_fail = len(scatter_data[(scatter_data[subject_x] < 10) & (scatter_data[subject_y] < 10)])
            total = len(scatter_data)
            
            col1, col2, col3, col4 = st.columns(4)
            with col1:
                st.metric("✅ ناجحون في الاثنين", f"{both_pass} ({both_pass/total*100:.0f}%)")
            with col2:
                st.metric(f"📗 ناجحون في {subject_x[:10]}.. فقط", f"{x_only} ({x_only/total*100:.0f}%)")
            with col3:
                st.metric(f"📘 ناجحون في {subject_y[:10]}.. فقط", f"{y_only} ({y_only/total*100:.0f}%)")
            with col4:
                st.metric("❌ راسبون في الاثنين", f"{both_fail} ({both_fail/total*100:.0f}%)")

    # Insights
    st.markdown("### 💡 استنتاجات تحليل الارتباط")
    
    # Find the most correlated pair
    if len(corr_df) > 0:
        strongest = corr_df.iloc[0]
        weakest = corr_df.iloc[-1]
        
        avg_correlation = corr_df['معامل الارتباط'].mean()
        
        if avg_correlation >= 0.5:
            st.success(f"🎯 **ترابط عام قوي:** متوسط الارتباط بين المواد هو {avg_correlation:.2f}. هذا يشير إلى أن التلاميذ المتفوقين يميلون للتفوق في معظم المواد.")
        elif avg_correlation >= 0.3:
            st.info(f"📊 **ترابط متوسط:** متوسط الارتباط {avg_correlation:.2f}. بعض المواد مترابطة والبعض الآخر مستقل.")
        else:
            st.warning(f"⚠️ **ترابط ضعيف:** متوسط الارتباط {avg_correlation:.2f}. كل مادة تتطلب مهارات مختلفة.")
        
        st.caption(f"🔗 أقوى علاقة: {strongest['المادة 1']} ↔ {strongest['المادة 2']} ({strongest['معامل الارتباط']:.2f})")
        st.caption(f"⛓️ أضعف علاقة: {weakest['المادة 1']} ↔ {weakest['المادة 2']} ({weakest['معامل الارتباط']:.2f})")

else:
    st.warning("⚠️ لا توجد بيانات كافية لحساب الارتباطات. يجب توفر بيانات 5 تلاميذ على الأقل.")

st.markdown("---")

# Individual Gap Analysis - At-Risk Report
st.header("🚨 تحليل الفجوات الفردية - تقرير التلاميذ المعرضين للخطر")

st.markdown("""
**تحديد التلاميذ الذين يحتاجون تدخلاً:** تحليل شامل للتلاميذ على حافة النجاح، والمتميزين، والذين يعانون من ضعف في مواد معينة.
""")

if 'المعدل' in df_filtered.columns:
    # Calculate statistics for classification
    avg_mean = df_filtered['المعدل'].dropna().mean()
    avg_std = df_filtered['المعدل'].dropna().std()
    
    # Classify students
    df_analysis = df_filtered[['ر.ت', 'رقم التلميذ', 'اسم التلميذ', 'المعدل'] + [col for col in subject_columns if col != 'المعدل' and col in df_filtered.columns]].copy()
    df_analysis = df_analysis.dropna(subset=['المعدل'])
    
    # Categories
    borderline_low = df_analysis[(df_analysis['المعدل'] >= 9) & (df_analysis['المعدل'] < 10)]
    borderline_high = df_analysis[(df_analysis['المعدل'] >= 10) & (df_analysis['المعدل'] < 11)]
    at_risk = df_analysis[df_analysis['المعدل'] < 9]
    excellent = df_analysis[df_analysis['المعدل'] >= avg_mean + 1.5 * avg_std]
    outliers_top = df_analysis[df_analysis['المعدل'] >= avg_mean + 2 * avg_std]
    
    # Summary metrics
    col1, col2, col3, col4 = st.columns(4)
    
    with col1:
        st.metric(
            "🔴 معرضون لخطر الهضر المدرسي",
            len(at_risk),
            help="تلاميذ معدلهم أقل من 9 - يحتاجون تدخلاً عاجلاً"
        )
    
    with col2:
        st.metric(
            "🟡 على الحافة",
            len(borderline_low),
            help="تلاميذ معدلهم بين 9 و 10 - قريبون من الرسوب"
        )
    
    with col3:
        st.metric(
            "🟢 ناجحون بصعوبة",
            len(borderline_high),
            help="تلاميذ معدلهم بين 10 و 11 - نجحوا لكن يحتاجون دعماً"
        )
    
    with col4:
        st.metric(
            "⭐ متميزون",
            len(excellent),
            help=f"تلاميذ معدلهم أعلى من {avg_mean + 1.5 * avg_std:.2f}"
        )
    
    # Tab layout for different categories
    tab1, tab2, tab3, tab4 = st.tabs(["🔴 المعرضون لخطر الهضر المدرسي", "🟡 على الحافة", "⭐ المتميزون", "📊 تحليل الضعف"])
    
    with tab1:
        st.markdown("### 🔴 التلاميذ المعرضون لخطر الهضر المدرسي (معدل < 9)")
        if len(at_risk) > 0:
            st.warning(f"⚠️ يوجد **{len(at_risk)}** تلاميذ بحاجة إلى تدخل عاجل!")
            
            for idx, row in at_risk.iterrows():
                with st.expander(f"📋 {row['اسم التلميذ']} - المعدل: {row['المعدل']:.2f}"):
                    # Find weakest subjects
                    subject_scores = {}
                    for col in subject_columns:
                        if col != 'المعدل' and col in df_filtered.columns and pd.notna(row.get(col)):
                            subject_scores[col] = row[col]
                    
                    if subject_scores:
                        sorted_subjects = sorted(subject_scores.items(), key=lambda x: x[1])
                        
                        st.markdown("**🔻 أضعف المواد (تحتاج تدخلاً):**")
                        for subj, score in sorted_subjects[:3]:
                            color = "red" if score < 10 else "green"
                            gap = 10 - score
                            st.markdown(f"- **{subj}**: :red[{score:.2f}] (يحتاج +{gap:.2f} للنجاح)")
                        
                        # Calculate what's needed
                        current_avg = row['المعدل']
                        points_needed = (10 - current_avg) * len(subject_scores)
                        st.info(f"💡 يحتاج إلى رفع مجموع نقاطه بـ **{points_needed:.1f}** نقطة للوصول للمعدل 10")
        else:
            st.success("✅ لا يوجد تلاميذ معرضون لخطر الهضر المدرسي!")
    
    with tab2:
        st.markdown("### 🟡 التلاميذ على الحافة (معدل 9-10)")
        if len(borderline_low) > 0:
            st.info(f"📊 يوجد **{len(borderline_low)}** تلاميذ قريبون جداً من خط النجاح")
            
            for idx, row in borderline_low.iterrows():
                with st.expander(f"📋 {row['اسم التلميذ']} - المعدل: {row['المعدل']:.2f}"):
                    subject_scores = {}
                    for col in subject_columns:
                        if col != 'المعدل' and col in df_filtered.columns and pd.notna(row.get(col)):
                            subject_scores[col] = row[col]
                    
                    if subject_scores:
                        sorted_subjects = sorted(subject_scores.items(), key=lambda x: x[1])
                        failing_subjects = [(s, sc) for s, sc in sorted_subjects if sc < 10]
                        
                        if failing_subjects:
                            st.markdown("**🎯 المواد التي تسحب المعدل للأسفل:**")
                            for subj, score in failing_subjects[:3]:
                                gap = 10 - score
                                st.markdown(f"- **{subj}**: :red[{score:.2f}] (فجوة: {gap:.2f})")
                            
                            # Quick win suggestion
                            easiest_fix = failing_subjects[0]
                            st.success(f"💡 **أسهل تحسين:** رفع درجة **{easiest_fix[0]}** من {easiest_fix[1]:.2f} إلى 10 سيرفع المعدل بشكل ملحوظ")
                        else:
                            st.success("جميع المواد فوق 10 - المعدل منخفض بسبب بعض الدرجات القريبة من 10")
        else:
            st.success("✅ لا يوجد تلاميذ على حافة الرسوب!")
        
        # Also show borderline successful students
        st.markdown("### 🟢 ناجحون لكن يحتاجون دعماً (معدل 10-11)")
        if len(borderline_high) > 0:
            st.info(f"📊 يوجد **{len(borderline_high)}** تلاميذ نجحوا بفارق بسيط")
            
            borderline_high_sorted = borderline_high.sort_values('المعدل')
            for idx, row in borderline_high_sorted.head(5).iterrows():
                subject_scores = {col: row[col] for col in subject_columns 
                                if col != 'المعدل' and col in df_filtered.columns and pd.notna(row.get(col))}
                if subject_scores:
                    weakest = min(subject_scores.items(), key=lambda x: x[1])
                    st.caption(f"• {row['اسم التلميذ']} ({row['المعدل']:.2f}) - أضعف مادة: {weakest[0]} ({weakest[1]:.2f})")
    
    with tab3:
        st.markdown("### ⭐ التلاميذ المتميزون - نموذج التفوق")
        
        if len(excellent) > 0:
            st.success(f"🌟 يوجد **{len(excellent)}** تلاميذ متميزين يمكن اعتبارهم نموذجاً!")
            
            # Top performers
            top_students = excellent.nlargest(5, 'المعدل')
            
            for idx, row in top_students.iterrows():
                with st.expander(f"🏆 {row['اسم التلميذ']} - المعدل: {row['المعدل']:.2f}", expanded=True):
                    subject_scores = {}
                    for col in subject_columns:
                        if col != 'المعدل' and col in df_filtered.columns and pd.notna(row.get(col)):
                            subject_scores[col] = row[col]
                    
                    if subject_scores:
                        sorted_subjects = sorted(subject_scores.items(), key=lambda x: x[1], reverse=True)
                        
                        col1, col2 = st.columns(2)
                        with col1:
                            st.markdown("**💪 أقوى المواد:**")
                            for subj, score in sorted_subjects[:3]:
                                st.markdown(f"- **{subj}**: :green[{score:.2f}]")
                        
                        with col2:
                            st.markdown("**📈 مجال للتحسين:**")
                            for subj, score in sorted_subjects[-2:]:
                                st.markdown(f"- **{subj}**: {score:.2f}")
            
            # Outlier analysis
            if len(outliers_top) > 0:
                st.markdown("### 🚀 التلاميذ الاستثنائيون (Outliers)")
                st.info(f"هؤلاء التلاميذ ({len(outliers_top)}) يتفوقون بشكل استثنائي على زملائهم")
                
                for idx, row in outliers_top.iterrows():
                    gap_from_avg = row['المعدل'] - avg_mean
                    st.caption(f"🌟 **{row['اسم التلميذ']}**: {row['المعدل']:.2f} (+{gap_from_avg:.2f} عن المتوسط)")
        else:
            st.info("لا يوجد تلاميذ متميزون بشكل استثنائي في هذه المجموعة")
    
    with tab4:
        st.markdown("### 📊 تحليل نقاط الضعف حسب المادة")
        
        # Find subjects where most students struggle
        subject_failure_analysis = []
        for col in subject_columns:
            if col != 'المعدل' and col in df_filtered.columns:
                subject_data = df_filtered[col].dropna()
                if len(subject_data) > 0:
                    failing_count = (subject_data < 10).sum()
                    failing_pct = (subject_data < 10).mean() * 100
                    avg_score = subject_data.mean()
                    subject_failure_analysis.append({
                        'المادة': col,
                        'عدد الراسبين': failing_count,
                        'نسبة الرسوب %': failing_pct,
                        'المتوسط': avg_score
                    })
        
        if subject_failure_analysis:
            failure_df = pd.DataFrame(subject_failure_analysis)
            failure_df = failure_df.sort_values('نسبة الرسوب %', ascending=False)
            
            # Visualization
            fig = px.bar(
                failure_df,
                x='المادة',
                y='نسبة الرسوب %',
                color='نسبة الرسوب %',
                color_continuous_scale='RdYlGn_r',
                text='عدد الراسبين'
            )
            fig.update_traces(texttemplate='%{text} تلميذ', textposition='outside')
            fig.update_layout(height=400, title="نسبة الرسوب في كل مادة")
            fig.add_hline(y=50, line_dash="dash", line_color="red", annotation_text="خط الخطر (50%)")
            st.plotly_chart(fig, use_container_width=True)
            
            # Critical subjects
            critical_subjects = failure_df[failure_df['نسبة الرسوب %'] > 50]
            if len(critical_subjects) > 0:
                st.error(f"⚠️ **مواد حرجة** (أكثر من 50% رسوب): {', '.join(critical_subjects['المادة'].tolist())}")
            
            # Students who fail in multiple subjects
            st.markdown("### 📉 التلاميذ الذين يرسبون في عدة مواد")
            
            multi_fail_students = []
            for idx, row in df_analysis.iterrows():
                failing_subjects = []
                for col in subject_columns:
                    if col != 'المعدل' and col in df_filtered.columns and pd.notna(row.get(col)):
                        if row[col] < 10:
                            failing_subjects.append(col)
                
                if len(failing_subjects) >= 3:
                    multi_fail_students.append({
                        'التلميذ': row['اسم التلميذ'],
                        'المعدل': row['المعدل'],
                        'عدد المواد الراسب فيها': len(failing_subjects),
                        'المواد': ', '.join(failing_subjects[:5])
                    })
            
            if multi_fail_students:
                multi_fail_df = pd.DataFrame(multi_fail_students)
                multi_fail_df = multi_fail_df.sort_values('عدد المواد الراسب فيها', ascending=False)
                
                st.dataframe(multi_fail_df, use_container_width=True, hide_index=True)
                
                worst_case = multi_fail_df.iloc[0]
                st.warning(f"⚠️ الحالة الأكثر خطورة: **{worst_case['التلميذ']}** يرسب في **{worst_case['عدد المواد الراسب فيها']}** مواد")
            else:
                st.success("✅ لا يوجد تلاميذ يرسبون في 3 مواد أو أكثر")

    # Recommendations
    st.markdown("### 💡 توصيات للتدخل")
    
    recommendations = []
    
    if len(at_risk) > 0:
        recommendations.append(f"🔴 **تدخل عاجل:** {len(at_risk)} تلاميذ يحتاجون دعماً مكثفاً فورياً")
    
    if len(borderline_low) > 0:
        recommendations.append(f"🟡 **متابعة دقيقة:** {len(borderline_low)} تلاميذ على حافة الرسوب يحتاجون دعماً مستهدفاً")
    
    if len(critical_subjects) > 0 if 'critical_subjects' in dir() else False:
        recommendations.append(f"📚 **مراجعة طرق التدريس:** المواد الحرجة تحتاج اهتماماً خاصاً")
    
    if len(excellent) > 0:
        recommendations.append(f"⭐ **برنامج تميز:** {len(excellent)} تلاميذ متميزين يمكن إشراكهم في مساعدة زملائهم")
    
    for rec in recommendations:
        st.markdown(f"- {rec}")

else:
    st.warning("⚠️ لا يوجد عمود 'المعدل' في البيانات")

st.markdown("---")

# Raw Data Table
st.header("📋 جميع بيانات التلاميذ")
# Dynamically select available columns to avoid KeyError
display_cols = ['ر.ت', 'رقم التلميذ', 'اسم التلميذ'] + [col for col in subject_columns if col in df_filtered.columns]
# Remove duplicates if any (e.g. if 'المعدل' is in both lists)
display_cols = list(dict.fromkeys(display_cols))
st.dataframe(df_filtered[display_cols], 
             use_container_width=True, height=400)

# Download option
st.markdown("---")

col_csv, col_ppt = st.columns(2)

with col_csv:
    # Add UTF-8 BOM for Excel to recognize Arabic characters
    csv = '\ufeff' + df_filtered.to_csv(index=False)
    st.download_button(
        label="📥 تحميل البيانات كـ CSV",
        data=csv.encode('utf-8'),
        file_name=f"student_data_statistics.csv",
        mime="text/csv"
    )

with col_ppt:
    st.subheader("📊 إنشاء عرض تقديمي")
    
    # Get all available classes
    all_classes = list(df['الفصل'].unique())
    
    # Option to combine all classes
    combine_all_classes = st.checkbox(
        "دمج جميع الفصول في عرض واحد",
        value=True,
        help="عند التفعيل، سيتم دمج بيانات جميع الفصول المختارة في إحصائيات موحدة"
    )
    
    # Multi-select for classes to include in presentation
    selected_classes_ppt = st.multiselect(
        "اختر الفصول للعرض التقديمي:",
        options=all_classes,
        default=all_classes,
        help="اختر الفصول التي تريد تضمينها في العرض التقديمي"
    )
    
    if len(selected_classes_ppt) == 0:
        st.warning("⚠️ الرجاء اختيار فصل واحد على الأقل")
    
    # Filter data for presentation based on selected classes
    df_ppt = df[df['الفصل'].isin(selected_classes_ppt)].copy()
    df_ppt = df_ppt.dropna(subset=['اسم التلميذ'])
    
    # Show summary of selection
    if st.button("📊 إنشاء العرض التقديمي (PPTX)", disabled=len(selected_classes_ppt) == 0):
        with st.spinner("جاري إنشاء العرض التقديمي..."):
            try:
                # Initialize presentation
                prs = Presentation()
                
                # Set 16:9 widescreen layout
                prs.slide_width = Inches(13.333)
                prs.slide_height = Inches(7.5)
                
                # Create slides based on selection
                if combine_all_classes:
                    generate_slides_for_data(prs, df_ppt, subject_columns, selected_classes_ppt)
                else:
                    for class_name in selected_classes_ppt:
                        class_df = df_ppt[df_ppt['الفصل'] == class_name]
                        if len(class_df) > 0:
                            generate_slides_for_data(prs, class_df, subject_columns, [class_name], title_suffix=f"- {class_name}")
                
                # Save to buffer
                ppt_buffer = io.BytesIO()
                prs.save(ppt_buffer)
                ppt_buffer.seek(0)
                
                st.success("✅ تم إنشاء العرض التقديمي بنجاح!")
                st.download_button(
                    label="📥 تحميل العرض التقديمي",
                    data=ppt_buffer,
                    file_name="student_analysis_report.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
            except Exception as e:
                st.error(f"❌ حدث خطأ أثناء إنشاء العرض التقديمي: {str(e)}")
                import traceback
                st.code(traceback.format_exc())
    
